from __future__ import annotations

import re
from datetime import datetime, timezone
from io import BytesIO
from typing import Any, Awaitable, Callable

import aiosqlite
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.config import Settings
from app.db_paths import common_db_path
from app.schemas import (
    ChartBarRequestDto,
    ChartPptExportRequestDto,
    ChartStackedMatrixRowDto,
    ChartStackedRequestDto,
    ChartStackedResolvedVersionDto,
    ChartStackedResponseDto,
    ChartStackedSeriesDto,
    ChartVersionItemDto,
)


def _extract_report_code(level_value: Any) -> str | None:
    text = str(level_value or "").strip()
    if not text:
        return None
    m = re.match(r"^([A-Z]\d+)\b", text)
    return m.group(1) if m else None


def _stacked_period_labels(granularity: str) -> list[str]:
    if granularity == "quarter":
        return ["Q1", "Q2", "Q3", "Q4"]
    return [f"M{i:02d}" for i in range(1, 13)]


def build_chart_write_router(
    *,
    settings: Settings,
    chart_version_options_provider: Callable[[], Awaitable[list[ChartVersionItemDto]]],
    extract_data_acct_code_from_name: Callable[[str], str | None],
) -> APIRouter:
    router = APIRouter()

    @router.post("/api/chart/stacked", response_model=ChartStackedResponseDto)
    async def chart_stacked(req: ChartStackedRequestDto):
        selected_versions = req.selected_versions
        if not req.use_all_versions and not selected_versions:
            raise HTTPException(status_code=400, detail="请选择至少一个展示版本")
        if len(selected_versions) > 5:
            raise HTTPException(status_code=400, detail="最多选择5个展示版本")

        options = await chart_version_options_provider()
        option_map: dict[tuple[int, int], ChartVersionItemDto] = {
            (opt.data_file_id, opt.version_id): opt for opt in options
        }
        if req.use_all_versions:
            resolved_options = options
        else:
            resolved_options = []
            for sel in selected_versions:
                key = (int(sel.data_file_id), int(sel.version_id))
                opt = option_map.get(key)
                if opt is not None:
                    resolved_options.append(opt)
        if not resolved_options:
            return ChartStackedResponseDto(note="没有可用版本数据")

        async with aiosqlite.connect(common_db_path()) as cdb:
            await cdb.execute("PRAGMA foreign_keys = ON")
            cur_reports = await cdb.execute(
                """
                SELECT report_acct_code, report_acct_name, parent_code, is_summary
                FROM report_account
                """
            )
            report_rows = await cur_reports.fetchall()
            cur_mapping = await cdb.execute(
                "SELECT report_acct_code, data_acct_code FROM report_data_mapping"
            )
            mapping_rows = await cur_mapping.fetchall()
            cur_data = await cdb.execute(
                "SELECT data_acct_code, data_acct_name, value_type FROM data_account"
            )
            data_rows = await cur_data.fetchall()

        selected_report_code = req.report_acct_code.strip().upper()
        report_name_map: dict[str, str] = {}
        children_map: dict[str, list[str]] = {}
        summary_map: dict[str, bool] = {}
        for row in report_rows:
            code = str(row[0])
            report_name_map[code] = str(row[1])
            parent_code = str(row[2]) if row[2] else None
            summary_map[code] = bool(row[3])
            if parent_code:
                children_map.setdefault(parent_code, []).append(code)
        if selected_report_code not in report_name_map:
            raise HTTPException(status_code=400, detail=f"报告科目 {selected_report_code} 不存在")

        def _collect_descendants(root_code: str) -> set[str]:
            stack = [root_code]
            out: set[str] = set()
            while stack:
                code = stack.pop()
                if code in out:
                    continue
                out.add(code)
                stack.extend(children_map.get(code, []))
            return out

        mapping_by_report: dict[str, set[str]] = {}
        for row in mapping_rows:
            rcode = str(row[0])
            dcode = str(row[1])
            mapping_by_report.setdefault(rcode, set()).add(dcode)
        data_name_map: dict[str, str] = {str(row[0]): str(row[1]) for row in data_rows}
        data_value_type_map: dict[str, str] = {str(row[0]): str(row[2] or "") for row in data_rows}

        child_codes = sorted(children_map.get(selected_report_code, []))
        direct_data_codes = sorted(mapping_by_report.get(selected_report_code, set()))
        if child_codes:
            segment_roots = child_codes
        elif summary_map.get(selected_report_code, False) and direct_data_codes:
            segment_roots = [f"DATA::{dcode}" for dcode in direct_data_codes]
        else:
            segment_roots = [selected_report_code]
        selected_descendants = _collect_descendants(selected_report_code)

        segment_data_codes: dict[str, set[str]] = {}
        segment_label_map: dict[str, str] = {}
        for segment in segment_roots:
            if segment.startswith("DATA::"):
                dcode = segment.split("::", 1)[1]
                segment_data_codes[segment] = {dcode}
                segment_label_map[segment] = f"{dcode} {data_name_map.get(dcode, '')}".strip()
            else:
                descendants = _collect_descendants(segment)
                dset: set[str] = set()
                for dreport in descendants:
                    dset.update(mapping_by_report.get(dreport, set()))
                segment_data_codes[segment] = dset
                segment_label_map[segment] = f"{segment} {report_name_map.get(segment, '')}".strip()

        if child_codes:
            covered_codes: set[str] = set()
            for segment in child_codes:
                covered_codes.update(segment_data_codes.get(segment, set()))
            for dcode in direct_data_codes:
                if dcode in covered_codes:
                    continue
                seg_key = f"DATA::{dcode}"
                segment_data_codes[seg_key] = {dcode}
                segment_label_map[seg_key] = f"{dcode} {data_name_map.get(dcode, '')}".strip()

        # 去重：同一 data_acct_code 只归属一个 segment，避免重复汇总。
        segment_order = sorted(segment_roots)
        if child_codes:
            segment_order.extend(
                sorted([key for key in segment_data_codes.keys() if key.startswith("DATA::")])
            )
        code_to_segment: dict[str, str] = {}
        for segment in segment_order:
            for dcode in sorted(segment_data_codes.get(segment, set())):
                if dcode not in code_to_segment:
                    code_to_segment[dcode] = segment

        segment_value_type_map: dict[str, str | None] = {}
        for segment in segment_order:
            value_types = {
                data_value_type_map.get(dcode, "")
                for dcode in segment_data_codes.get(segment, set())
                if data_value_type_map.get(dcode, "")
            }
            if len(value_types) == 1:
                segment_value_type_map[segment] = next(iter(value_types))
            else:
                segment_value_type_map[segment] = None

        effective_options = sorted(
            resolved_options,
            key=lambda o: (o.year, o.version_id),
        )
        single_version = len(effective_options) == 1
        categories = (
            _stacked_period_labels(req.single_version_granularity)
            if single_version
            else [f"{opt.year}-V{opt.version_id} {opt.version_name}" for opt in effective_options]
        )
        category_index_map = {name: idx for idx, name in enumerate(categories)}

        series_abs_map: dict[str, list[float]] = {
            seg: [0.0] * len(categories) for seg in segment_order
        }

        for x_idx, opt in enumerate(effective_options):
            budget_path = settings.data_dir / opt.data_file_name
            if not budget_path.exists():
                continue
            async with aiosqlite.connect(budget_path) as bdb:
                await bdb.execute("PRAGMA foreign_keys = ON")
                cur = await bdb.execute(
                    """
                    SELECT data_code_name, month, quarter, value,
                           report_level1, report_level2, report_level3, report_level4, report_level5
                    FROM budget_summary
                    WHERE version_id = ?
                    """,
                    (opt.version_id,),
                )
                rows = await cur.fetchall()

            for row in rows:
                data_code = extract_data_acct_code_from_name(str(row[0] or ""))
                if not data_code:
                    continue
                segment = code_to_segment.get(data_code)
                if not segment:
                    continue
                # 双重保险：确认该行 report path 属于所选科目子树
                level_codes = {_extract_report_code(row[idx]) for idx in range(4, 9)}
                level_codes.discard(None)
                if selected_report_code not in level_codes and not any(
                    code in level_codes for code in selected_descendants
                ):
                    continue
                value = float(row[3] or 0.0)
                if single_version:
                    period_key = str(row[2] if req.single_version_granularity == "quarter" else row[1])
                    c_idx = category_index_map.get(period_key)
                    if c_idx is None:
                        continue
                else:
                    c_idx = x_idx
                series_abs_map[segment][c_idx] += value

        series: list[ChartStackedSeriesDto] = []
        matrix_rows: list[ChartStackedMatrixRowDto] = []
        if req.stack_mode == "percent":
            totals = [0.0] * len(categories)
            for values in series_abs_map.values():
                for i, v in enumerate(values):
                    totals[i] += v
            for seg in segment_order:
                raw_values = series_abs_map[seg]
                percent_values = [
                    (raw_values[i] / totals[i] * 100.0) if totals[i] else 0.0
                    for i in range(len(categories))
                ]
                label = segment_label_map.get(seg, seg)
                series.append(ChartStackedSeriesDto(key=seg, label=label, values=percent_values))
                matrix_rows.append(ChartStackedMatrixRowDto(row_label=label, values=percent_values))
        else:
            for seg in segment_order:
                values = series_abs_map[seg]
                label = segment_label_map.get(seg, seg)
                series.append(ChartStackedSeriesDto(key=seg, label=label, values=values))
                matrix_rows.append(ChartStackedMatrixRowDto(row_label=label, values=values))

        resolved_versions = [
            ChartStackedResolvedVersionDto(
                data_file_id=opt.data_file_id,
                year=opt.year,
                version_id=opt.version_id,
                version_name=opt.version_name,
            )
            for opt in effective_options
        ]
        return ChartStackedResponseDto(
            categories=categories,
            series=series,
            matrix_headers=categories,
            matrix_rows=matrix_rows,
            resolved_versions=resolved_versions,
            note="stack_mode=percent 时数值单位为百分比",
        )

    @router.post("/api/chart/bar", response_model=ChartStackedResponseDto)
    async def chart_bar(req: ChartBarRequestDto):
        """柱状图：横轴为期间（单版本）或版本（多版本）；本科目合计为一组柱，下级科目为分组多柱。"""
        selected_versions = req.selected_versions
        if not req.use_all_versions and not selected_versions:
            raise HTTPException(status_code=400, detail="请选择至少一个展示版本")
        if len(selected_versions) > 5:
            raise HTTPException(status_code=400, detail="最多选择5个展示版本")

        scope = req.bar_compare_scope.strip().lower()

        options = await chart_version_options_provider()
        option_map: dict[tuple[int, int], ChartVersionItemDto] = {
            (opt.data_file_id, opt.version_id): opt for opt in options
        }
        if req.use_all_versions:
            resolved_options = options
        else:
            resolved_options = []
            for sel in selected_versions:
                key = (int(sel.data_file_id), int(sel.version_id))
                opt = option_map.get(key)
                if opt is not None:
                    resolved_options.append(opt)
        if not resolved_options:
            return ChartStackedResponseDto(note="没有可用版本数据")

        async with aiosqlite.connect(common_db_path()) as cdb:
            await cdb.execute("PRAGMA foreign_keys = ON")
            cur_reports = await cdb.execute(
                """
                SELECT report_acct_code, report_acct_name, parent_code, is_summary
                FROM report_account
                """
            )
            report_rows = await cur_reports.fetchall()
            cur_mapping = await cdb.execute(
                "SELECT report_acct_code, data_acct_code FROM report_data_mapping"
            )
            mapping_rows = await cur_mapping.fetchall()
            cur_data = await cdb.execute(
                "SELECT data_acct_code, data_acct_name, value_type FROM data_account"
            )
            data_rows = await cur_data.fetchall()

        selected_report_code = req.report_acct_code.strip().upper()
        report_name_map: dict[str, str] = {}
        children_map: dict[str, list[str]] = {}
        for row in report_rows:
            code = str(row[0])
            report_name_map[code] = str(row[1])
            parent_code = str(row[2]) if row[2] else None
            if parent_code:
                children_map.setdefault(parent_code, []).append(code)
        if selected_report_code not in report_name_map:
            raise HTTPException(status_code=400, detail=f"报告科目 {selected_report_code} 不存在")

        def _collect_descendants(root_code: str) -> set[str]:
            stack = [root_code]
            out: set[str] = set()
            while stack:
                code = stack.pop()
                if code in out:
                    continue
                out.add(code)
                stack.extend(children_map.get(code, []))
            return out

        mapping_by_report: dict[str, set[str]] = {}
        for row in mapping_rows:
            rcode = str(row[0])
            dcode = str(row[1])
            mapping_by_report.setdefault(rcode, set()).add(dcode)
        data_name_map: dict[str, str] = {str(row[0]): str(row[1]) for row in data_rows}
        data_value_type_map: dict[str, str] = {str(row[0]): str(row[2] or "") for row in data_rows}

        selected_descendants = _collect_descendants(selected_report_code)
        direct_data_codes = sorted(mapping_by_report.get(selected_report_code, set()))

        if scope == "self":
            segment_order = [selected_report_code]
            dset: set[str] = set()
            for dreport in selected_descendants:
                dset.update(mapping_by_report.get(dreport, set()))
            segment_data_codes = {selected_report_code: dset}
            segment_label_map = {
                selected_report_code: f"{selected_report_code} {report_name_map.get(selected_report_code, '')}".strip()
            }
        else:
            child_codes = sorted(children_map.get(selected_report_code, []))
            if not child_codes and not direct_data_codes:
                raise HTTPException(
                    status_code=400,
                    detail="所选报告科目没有下级科目，请改用本科目比较",
                )
            segment_order = child_codes[:]
            segment_data_codes = {}
            segment_label_map = {}
            for segment in segment_order:
                descendants = _collect_descendants(segment)
                seg_d: set[str] = set()
                for dreport in descendants:
                    seg_d.update(mapping_by_report.get(dreport, set()))
                segment_data_codes[segment] = seg_d
                segment_label_map[segment] = f"{segment} {report_name_map.get(segment, '')}".strip()
            covered_codes: set[str] = set()
            for segment in child_codes:
                covered_codes.update(segment_data_codes.get(segment, set()))
            for dcode in direct_data_codes:
                if dcode in covered_codes:
                    continue
                seg_key = f"DATA::{dcode}"
                segment_order.append(seg_key)
                segment_data_codes[seg_key] = {dcode}
                segment_label_map[seg_key] = f"{dcode} {data_name_map.get(dcode, '')}".strip()

        code_to_segment: dict[str, str] = {}
        for segment in segment_order:
            for dcode in sorted(segment_data_codes.get(segment, set())):
                if dcode not in code_to_segment:
                    code_to_segment[dcode] = segment

        segment_value_type_map: dict[str, str | None] = {}
        for segment in segment_order:
            value_types = {
                data_value_type_map.get(dcode, "")
                for dcode in segment_data_codes.get(segment, set())
                if data_value_type_map.get(dcode, "")
            }
            segment_value_type_map[segment] = next(iter(value_types)) if len(value_types) == 1 else None

        effective_options = sorted(
            resolved_options,
            key=lambda o: (o.year, o.version_id),
        )
        single_version = len(effective_options) == 1
        categories = (
            _stacked_period_labels(req.single_version_granularity)
            if single_version
            else [f"{opt.year}-V{opt.version_id} {opt.version_name}" for opt in effective_options]
        )
        category_index_map = {name: idx for idx, name in enumerate(categories)}

        series_abs_map: dict[str, list[float]] = {
            seg: [0.0] * len(categories) for seg in segment_order
        }

        for x_idx, opt in enumerate(effective_options):
            budget_path = settings.data_dir / opt.data_file_name
            if not budget_path.exists():
                continue
            async with aiosqlite.connect(budget_path) as bdb:
                await bdb.execute("PRAGMA foreign_keys = ON")
                cur = await bdb.execute(
                    """
                    SELECT data_code_name, month, quarter, value,
                           report_level1, report_level2, report_level3, report_level4, report_level5
                    FROM budget_summary
                    WHERE version_id = ?
                    """,
                    (opt.version_id,),
                )
                rows = await cur.fetchall()

            for row in rows:
                data_code = extract_data_acct_code_from_name(str(row[0] or ""))
                if not data_code:
                    continue
                segment = code_to_segment.get(data_code)
                if not segment:
                    continue
                # DATA:: 前缀表示“直接挂在选中科目上的数据科目”，以映射为准，不再依赖 report_level 路径。
                if not segment.startswith("DATA::"):
                    level_codes = {_extract_report_code(row[idx]) for idx in range(4, 9)}
                    level_codes.discard(None)
                    if selected_report_code not in level_codes and not any(
                        code in level_codes for code in selected_descendants
                    ):
                        continue
                value = float(row[3] or 0.0)
                if single_version:
                    period_key = str(row[2] if req.single_version_granularity == "quarter" else row[1])
                    c_idx = category_index_map.get(period_key)
                    if c_idx is None:
                        continue
                else:
                    c_idx = x_idx
                series_abs_map[segment][c_idx] += value

        series: list[ChartStackedSeriesDto] = []
        matrix_rows: list[ChartStackedMatrixRowDto] = []
        for seg in segment_order:
            values = series_abs_map[seg]
            label = segment_label_map.get(seg, seg)
            value_type = segment_value_type_map.get(seg)
            series.append(
                ChartStackedSeriesDto(key=seg, label=label, values=values, value_type=value_type)
            )
            matrix_rows.append(
                ChartStackedMatrixRowDto(row_label=label, values=values, value_type=value_type)
            )

        resolved_versions = [
            ChartStackedResolvedVersionDto(
                data_file_id=opt.data_file_id,
                year=opt.year,
                version_id=opt.version_id,
                version_name=opt.version_name,
            )
            for opt in effective_options
        ]
        return ChartStackedResponseDto(
            categories=categories,
            series=series,
            matrix_headers=categories,
            matrix_rows=matrix_rows,
            resolved_versions=resolved_versions,
            note=None,
        )

    @router.post("/api/chart/export-ppt")
    async def chart_export_ppt(req: ChartPptExportRequestDto):
        categories = [str(c) for c in req.categories]
        if not categories:
            raise HTTPException(status_code=400, detail="categories 不能为空")
        if not req.series:
            raise HTTPException(status_code=400, detail="series 不能为空")

        for s in req.series:
            if len(s.values) != len(categories):
                raise HTTPException(
                    status_code=400,
                    detail=f"系列 {s.name} 的数据点数量与 categories 不一致",
                )

        prs = Presentation()
        # 固定为 16:9 页面，避免默认 4:3 页面导致图表/矩阵越界。
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        emu_per_inch = 914400
        slide_width_in = prs.slide_width / emu_per_inch
        slide_height_in = prs.slide_height / emu_per_inch
        left_margin_in = 0.45
        right_margin_in = 0.45
        top_margin_in = 0.18
        bottom_margin_in = 0.18
        content_width_in = slide_width_in - left_margin_in - right_margin_in

        title_box = slide.shapes.add_textbox(
            Inches(left_margin_in), Inches(top_margin_in), Inches(content_width_in), Inches(0.40)
        )
        title_box.text_frame.clear()
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = req.title
        title_para.font.size = Pt(16)
        title_para.font.bold = True

        subtitle_top_in = top_margin_in + 0.40 + 0.04
        chart_top_in = subtitle_top_in
        if req.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(left_margin_in), Inches(subtitle_top_in), Inches(content_width_in), Inches(0.24)
            )
            subtitle_box.text_frame.clear()
            subtitle_para = subtitle_box.text_frame.paragraphs[0]
            subtitle_para.text = req.subtitle
            subtitle_para.font.size = Pt(9)
            chart_top_in = subtitle_top_in + 0.24 + 0.08

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in req.series:
            chart_data.add_series(s.name, s.values)

        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }
        chart_type = chart_type_map[req.chart_type]

        has_matrix = bool(req.matrix_headers) and bool(req.matrix_rows)
        matrix_top_in = 0.0
        matrix_height_in = 0.0
        matrix_note_top_in = 0.0
        matrix_note_height_in = 0.0
        if has_matrix:
            matrix_note_height_in = 0.16
            matrix_height_in = min(2.40, slide_height_in * 0.33)
            matrix_note_top_in = slide_height_in - bottom_margin_in - matrix_note_height_in
            matrix_top_in = matrix_note_top_in - 0.03 - matrix_height_in
            chart_bottom_limit_in = matrix_top_in - 0.10
        else:
            chart_bottom_limit_in = slide_height_in - bottom_margin_in
        chart_height = max(2.60, chart_bottom_limit_in - chart_top_in)
        chart_shape = slide.shapes.add_chart(
            chart_type,
            Inches(left_margin_in),
            Inches(chart_top_in),
            Inches(content_width_in),
            Inches(chart_height),
            chart_data,
        )
        chart = chart_shape.chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
        chart.has_title = False

        if req.chart_type in {"pie", "doughnut"}:
            plot = chart.plots[0]
            plot.has_data_labels = True
            labels = plot.data_labels
            is_doughnut = req.chart_type == "doughnut"
            labels.show_category_name = True
            labels.show_percentage = not is_doughnut
            labels.show_value = is_doughnut
            labels.number_format = "0.00" if is_doughnut else "0.00%"
            labels.position = XL_LABEL_POSITION.BEST_FIT
            if hasattr(labels, "separator"):
                labels.separator = "\n"
            if hasattr(labels, "font"):
                labels.font.size = Pt(7)
            # python-pptx 在部分模板中不会稳定应用 plot 级字体，逐点设置兜底。
            for series in chart.series:
                for point in series.points:
                    point_label = point.data_label
                    point_label.show_percentage = not is_doughnut
                    point_label.show_value = is_doughnut
                    point_label.show_category_name = True
                    point_label.number_format = "0.00" if is_doughnut else "0.00%"
                    point_label.position = XL_LABEL_POSITION.BEST_FIT
                    point_label.font.size = Pt(7)
        else:
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(7)
            value_axis = chart.value_axis
            value_axis.tick_labels.font.size = Pt(7)

        # 若调用侧开启了数据标签，统一压小字号，避免标签过大遮挡图表。
        for plot in chart.plots:
            if getattr(plot, "has_data_labels", False):
                data_labels = plot.data_labels
                if hasattr(data_labels, "font"):
                    data_labels.font.size = Pt(7 if req.chart_type in {"pie", "doughnut"} else 7)

        # 在图表下方渲染数据矩阵，保持导出页信息完整。
        if has_matrix:
            max_data_columns = 8
            max_data_rows = 10

            matrix_headers = [str(h) for h in req.matrix_headers]
            matrix_rows: list[tuple[str, list[str]]] = []
            for row in req.matrix_rows:
                values = [str(v) for v in row.values]
                if len(values) < len(matrix_headers):
                    values = values + [""] * (len(matrix_headers) - len(values))
                elif len(values) > len(matrix_headers):
                    values = values[: len(matrix_headers)]
                matrix_rows.append((str(row.label), values))

            clipped_columns = len(matrix_headers) > max_data_columns
            clipped_rows = len(matrix_rows) > max_data_rows
            display_headers = matrix_headers[:max_data_columns]
            display_rows = matrix_rows[:max_data_rows]

            table_rows = 1 + len(display_rows)
            table_cols = 1 + len(display_headers)
            table_shape = slide.shapes.add_table(
                table_rows,
                table_cols,
                Inches(left_margin_in),
                Inches(matrix_top_in),
                Inches(content_width_in),
                Inches(matrix_height_in),
            )
            table = table_shape.table
            first_col_width = min(4.20, max(3.20, content_width_in * 0.34))
            table.columns[0].width = Inches(first_col_width)
            value_col_width = Inches(max(0.60, (content_width_in - first_col_width) / max(1, len(display_headers))))
            for col_idx in range(1, table_cols):
                table.columns[col_idx].width = value_col_width

            head_cell = table.cell(0, 0)
            head_cell.text = "维度项"
            head_para = head_cell.text_frame.paragraphs[0]
            head_para.font.bold = True
            head_para.font.size = Pt(11)
            head_para.alignment = PP_ALIGN.CENTER

            for col_idx, header in enumerate(display_headers, start=1):
                cell = table.cell(0, col_idx)
                cell.text = header
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.CENTER

            for row_idx, (row_label, values) in enumerate(display_rows, start=1):
                label_cell = table.cell(row_idx, 0)
                label_cell.text = row_label
                label_para = label_cell.text_frame.paragraphs[0]
                label_para.font.size = Pt(10)
                label_para.alignment = PP_ALIGN.LEFT

                for col_idx, value in enumerate(values[: len(display_headers)], start=1):
                    value_cell = table.cell(row_idx, col_idx)
                    value_cell.text = value
                    value_para = value_cell.text_frame.paragraphs[0]
                    value_para.font.size = Pt(10)
                    value_para.alignment = PP_ALIGN.RIGHT

            if clipped_columns or clipped_rows:
                note_box = slide.shapes.add_textbox(
                    Inches(left_margin_in), Inches(matrix_note_top_in), Inches(content_width_in), Inches(matrix_note_height_in)
                )
                note_text = "注：数据矩阵已按页面宽度/高度截取显示。"
                note_box.text_frame.text = note_text
                note_box.text_frame.paragraphs[0].font.size = Pt(7)

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        filename = f"pivot_chart_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.pptx"
        headers = {"Content-Disposition": f'attachment; filename="{filename}"'}
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers=headers,
        )

    return router
