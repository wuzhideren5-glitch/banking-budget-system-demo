"""PowerPoint template structure inspector for AI PPT binding workflows."""

from __future__ import annotations

import re
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.schemas import (
    SmartPptTemplateInspectResponse,
    SmartPptTemplateObjectRow,
    SmartPptTemplateSlideReportRow,
)

_WHITESPACE_RE = re.compile(r"\s+")


def _clean_excerpt(value: str | None, *, limit: int = 80) -> str | None:
    text = _WHITESPACE_RE.sub(" ", value or "").strip()
    if not text:
        return None
    return text if len(text) <= limit else f"{text[: limit - 1]}..."


def _safe_int(value: Any) -> int | None:
    try:
        return int(value)
    except Exception:
        return None


def _safe_len(value: Any) -> int | None:
    try:
        return len(value)
    except Exception:
        return None


class PptTemplateInspector:
    """Extract a stable, UI-friendly report from a PPTX template."""

    def inspect(self, template_path: Path) -> SmartPptTemplateInspectResponse:
        if not template_path.exists():
            raise FileNotFoundError(f"PPT template not found: {template_path}")

        prs = Presentation(str(template_path))
        slides: list[SmartPptTemplateSlideReportRow] = []
        warnings: list[str] = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_objects: list[SmartPptTemplateObjectRow] = []
            for path, shape in self._iter_shapes(slide.shapes):
                try:
                    slide_objects.append(self._inspect_shape(slide_index, path, shape))
                except Exception as exc:
                    warnings.append(f"第 {slide_index} 页形状 {path} 解析失败：{exc}")

            title = next((item.text_excerpt for item in slide_objects if item.text_excerpt), None)
            slides.append(
                SmartPptTemplateSlideReportRow(
                    slide_index=slide_index,
                    title=title,
                    object_count=len(slide_objects),
                    text_count=sum(1 for item in slide_objects if item.object_type == "text"),
                    table_count=sum(1 for item in slide_objects if item.object_type == "table"),
                    chart_count=sum(1 for item in slide_objects if item.object_type == "chart"),
                    picture_count=sum(1 for item in slide_objects if item.object_type == "picture"),
                    group_count=sum(1 for item in slide_objects if item.object_type == "group"),
                    other_count=sum(1 for item in slide_objects if item.object_type == "other"),
                    objects=slide_objects,
                )
            )

        return SmartPptTemplateInspectResponse(
            template_file_name=template_path.name,
            slide_count=len(slides),
            slide_width=int(prs.slide_width),
            slide_height=int(prs.slide_height),
            object_count=sum(slide.object_count for slide in slides),
            text_count=sum(slide.text_count for slide in slides),
            table_count=sum(slide.table_count for slide in slides),
            chart_count=sum(slide.chart_count for slide in slides),
            picture_count=sum(slide.picture_count for slide in slides),
            group_count=sum(slide.group_count for slide in slides),
            other_count=sum(slide.other_count for slide in slides),
            slides=slides,
            warnings=warnings,
        )

    def _iter_shapes(self, shapes: Iterable[Any], prefix: str = "") -> Iterable[tuple[str, Any]]:
        for index, shape in enumerate(shapes, start=1):
            path = f"{prefix}.{index}" if prefix else str(index)
            yield path, shape
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                yield from self._iter_shapes(shape.shapes, path)

    def _inspect_shape(self, slide_index: int, path: str, shape: Any) -> SmartPptTemplateObjectRow:
        object_type = self._shape_type(shape)
        row_count: int | None = None
        column_count: int | None = None
        chart_type: str | None = None
        text_excerpt: str | None = None

        if object_type == "table":
            table = shape.table
            row_count = _safe_len(table.rows)
            column_count = _safe_len(table.columns)
            text_excerpt = self._table_excerpt(table)
        elif object_type == "chart":
            chart_type = str(getattr(shape.chart, "chart_type", "") or "") or None
            text_excerpt = self._chart_title_excerpt(shape.chart)
        elif getattr(shape, "has_text_frame", False):
            text_excerpt = _clean_excerpt(getattr(shape, "text", None))

        return SmartPptTemplateObjectRow(
            object_id=f"s{slide_index}:{path}",
            shape_id=_safe_int(getattr(shape, "shape_id", None)),
            shape_name=str(getattr(shape, "name", "") or "") or None,
            object_type=object_type,
            text_excerpt=text_excerpt,
            chart_type=chart_type,
            row_count=row_count,
            column_count=column_count,
            left=_safe_int(getattr(shape, "left", None)),
            top=_safe_int(getattr(shape, "top", None)),
            width=_safe_int(getattr(shape, "width", None)),
            height=_safe_int(getattr(shape, "height", None)),
        )

    def _shape_type(self, shape: Any) -> str:
        if getattr(shape, "has_chart", False):
            return "chart"
        if getattr(shape, "has_table", False):
            return "table"
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            return "picture"
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            return "group"
        if getattr(shape, "has_text_frame", False) and _clean_excerpt(getattr(shape, "text", None)):
            return "text"
        return "other"

    def _table_excerpt(self, table: Any) -> str | None:
        cells: list[str] = []
        for row in list(table.rows)[:2]:
            for cell in list(row.cells)[:4]:
                excerpt = _clean_excerpt(getattr(cell, "text", None), limit=24)
                if excerpt:
                    cells.append(excerpt)
        return " / ".join(cells[:6]) or None

    def _chart_title_excerpt(self, chart: Any) -> str | None:
        try:
            if not getattr(chart, "has_title", False):
                return None
            return _clean_excerpt(chart.chart_title.text_frame.text)
        except Exception:
            return None
