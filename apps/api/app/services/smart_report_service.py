from __future__ import annotations

# ─── 模块辅助函数 ───
import json
import re
import shutil
import ast
import operator
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import aiosqlite
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
from fastapi import HTTPException, UploadFile

from app.db_bootstrap.derived_read_models import ensure_budget_summary_read_model_schema_async
from app.core.db_paths import budget_db_path, common_db_path
from app.services.runtime_metric_refs import load_confirmed_org_product_runtime_ref_codes
from app.schemas import (
    SmartReportAIBlock,
    SmartReportAIInspectionIssue,
    SmartReportAIInspectionResponse,
    SmartReportBlueprintDetail,
    SmartReportBlueprintGenerateResponse,
    SmartReportBlueprintPreviewResponse,
    SmartReportBlueprintRow,
    SmartReportBlueprintSaveRequest,
    SmartReportCalcMetricComponent,
    SmartReportCalcMetricRow,
    SmartReportCalcMetricUpsert,
    SmartReportGenerateRequest,
    SmartReportGenerateResponse,
    SmartReportInstanceRow,
    SmartReportPreviewRequest,
    SmartReportPreviewResponse,
    SmartReportTemplateCreateResponse,
    SmartReportTemplateRow,
    SmartReportTemplateVariableRow,
    SmartReportTemplateVariableUpsert,
)


PLACEHOLDER_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")
CHART_PLACEHOLDER_RE = re.compile(r"\{\{\s*chart\s*:\s*([^{}]+?)\s*\}\}", re.IGNORECASE)


def _iso_now() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _safe_code(raw: str) -> str:
    code = re.sub(r"[^A-Za-z0-9_-]+", "_", raw.strip())
    code = re.sub(r"_+", "_", code).strip("_")
    if not code:
        raise HTTPException(status_code=400, detail="模板编码不能为空")
    return code[:80]


def _json_loads_dict(raw: str | None) -> dict[str, Any]:
    if not raw:
        return {}
    try:
        value = json.loads(raw)
    except Exception:
        return {}
    return value if isinstance(value, dict) else {}


def _month_label(value: Any) -> str | None:
    if value is None or str(value).strip() == "":
        return None
    text = str(value).strip().upper()
    if re.fullmatch(r"M\d{2}", text):
        return text
    if re.fullmatch(r"\d{1,2}", text):
        num = int(text)
        if 1 <= num <= 12:
            return f"M{num:02d}"
    return text


def _year_label(value: Any) -> str:
    text = str(value or "").strip().upper()
    if re.fullmatch(r"Y\d{4}", text):
        return text
    if re.fullmatch(r"\d{4}", text):
        return f"Y{text}"
    return f"Y{datetime.now().year}"


def _plain_year(value: Any) -> int:
    label = _year_label(value)
    m = re.search(r"(\d{4})", label)
    return int(m.group(1)) if m else datetime.now().year


def _format_number(value: float, *, percentage: bool = False) -> str:
    if percentage:
        return f"{value * 100:.2f}%"
    abs_value = abs(value)
    if abs_value >= 10000:
        return f"{value:,.2f}"
    if value == int(value):
        return f"{int(value):,}"
    return f"{value:,.2f}"


_CALC_OPERATORS: dict[type[ast.AST], Any] = {
    ast.Add: operator.add,
    ast.Sub: operator.sub,
    ast.Mult: operator.mul,
    ast.Div: operator.truediv,
    ast.Pow: operator.pow,
    ast.USub: operator.neg,
    ast.UAdd: operator.pos,
}


# ─── 报告模板与蓝图目录 ───

class SmartReportService:
    def __init__(self, *, data_dir: Path, smart_ppt_service: Any = None, deepseek_client: Any = None) -> None:
        self.data_dir = data_dir
        self.template_dir = data_dir / "smart_report_templates"
        self.output_dir = data_dir / "smart_report_outputs"
        self.chart_cache_dir = data_dir / "chart_cache"
        self.smart_ppt_service = smart_ppt_service
        self.deepseek_client = deepseek_client

    async def list_templates(self) -> list[SmartReportTemplateRow]:
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                """
                SELECT t.template_id, t.template_code, t.template_name, t.template_type,
                       t.status, t.version_no, t.remark, t.created_at, t.updated_at,
                       COUNT(v.variable_id) AS variable_count
                FROM smart_report_template t
                LEFT JOIN smart_report_template_variable v ON v.template_id = t.template_id
                GROUP BY t.template_id
                ORDER BY t.updated_at DESC, t.template_id DESC
                """
            )
            rows = await cur.fetchall()
        return [
            SmartReportTemplateRow(
                template_id=int(r[0]),
                template_code=str(r[1]),
                template_name=str(r[2]),
                template_type=str(r[3]),
                status=str(r[4]),
                version_no=int(r[5]),
                remark=str(r[6]) if r[6] is not None else None,
                created_at=str(r[7]),
                updated_at=str(r[8]),
                variable_count=int(r[9] or 0),
            )
            for r in rows
        ]

    async def inspect_report_with_ai(self, file: UploadFile) -> SmartReportAIInspectionResponse:
        filename = file.filename or "未命名报告.docx"
        if not filename.lower().endswith(".docx"):
            raise HTTPException(status_code=400, detail="AI 报告解析第一版仅支持 .docx")
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="上传文件为空")
        tmp = self.template_dir / f"_ai_inspect_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.docx"
        self.template_dir.mkdir(parents=True, exist_ok=True)
        try:
            tmp.write_bytes(content)
            report_text = self._extract_docx_plain_text(tmp)
        finally:
            tmp.unlink(missing_ok=True)
        compact = re.sub(r"\s+", " ", report_text).strip()
        if not compact:
            raise HTTPException(status_code=400, detail="未从 Word 中提取到有效文本")
        catalog = await self._load_ai_binding_catalog()
        raw = self._call_ai_report_inspector(compact[:12000], catalog)
        parsed = self._parse_ai_inspection(raw) if raw else self._fallback_ai_inspection(compact)
        warnings = [] if raw else ["DeepSeek 未返回可用结果，已使用规则兜底解析。"]
        return SmartReportAIInspectionResponse(
            filename=filename,
            model=str(getattr(self.deepseek_client, "model", "") or ""),
            summary=str(parsed.get("summary") or ""),
            blocks=[
                SmartReportAIBlock(
                    block_id=str(item.get("block_id") or f"B{idx + 1}"),
                    block_type=str(item.get("block_type") or "text_block"),
                    text=str(item.get("text") or ""),
                    metrics=item.get("metrics") if isinstance(item.get("metrics"), list) else [],
                    analysis_rule_nl=str(item.get("analysis_rule_nl")) if item.get("analysis_rule_nl") else None,
                    structured_plan=item.get("structured_plan") if isinstance(item.get("structured_plan"), dict) else {},
                    confidence=float(item.get("confidence") or 0),
                )
                for idx, item in enumerate(parsed.get("blocks") or [])
                if isinstance(item, dict)
            ],
            issues=[
                SmartReportAIInspectionIssue(
                    issue_type=str(item.get("issue_type") or "needs_confirmation"),
                    text=str(item.get("text") or ""),
                    suggested_action=str(item.get("suggested_action") or ""),
                    candidates=item.get("candidates") if isinstance(item.get("candidates"), list) else [],
                    rule_preview=str(item.get("rule_preview")) if item.get("rule_preview") else None,
                )
                for item in (parsed.get("issues") or [])
                if isinstance(item, dict)
            ],
            assumptions=[str(x) for x in (parsed.get("assumptions") or [])],
            raw_text_excerpt=compact[:1200],
            warnings=warnings,
        )

    async def list_blueprints(self) -> list[SmartReportBlueprintRow]:
        async with aiosqlite.connect(common_db_path()) as db:
            cur = await db.execute(
                """
                SELECT blueprint_id, blueprint_name, source_filename, inspection_json, status,
                       output_file_path, last_generated_at, created_at, updated_at
                FROM smart_report_blueprint
                ORDER BY updated_at DESC, blueprint_id DESC
                LIMIT 100
                """
            )
            rows = await cur.fetchall()
        return [self._blueprint_row_from_db(row) for row in rows]

    async def save_blueprint(self, body: SmartReportBlueprintSaveRequest) -> SmartReportBlueprintDetail:
        now = _iso_now()
        inspection_json = json.dumps(body.inspection.model_dump(), ensure_ascii=False)
        async with aiosqlite.connect(common_db_path()) as db:
            cur = await db.execute(
                """
                INSERT INTO smart_report_blueprint (
                  blueprint_name, source_filename, inspection_json, status,
                  created_at, updated_at
                ) VALUES (?, ?, ?, 'draft', ?, ?)
                """,
                (body.blueprint_name.strip(), body.inspection.filename, inspection_json, now, now),
            )
            blueprint_id = int(cur.lastrowid)
            await db.commit()
        return await self.get_blueprint(blueprint_id)

    async def get_blueprint(self, blueprint_id: int) -> SmartReportBlueprintDetail:
        async with aiosqlite.connect(common_db_path()) as db:
            cur = await db.execute(
                """
                SELECT blueprint_id, blueprint_name, source_filename, inspection_json, status,
                       output_file_path, last_generated_at, created_at, updated_at
                FROM smart_report_blueprint
                WHERE blueprint_id = ?
                """,
                (blueprint_id,),
            )
            row = await cur.fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="报告蓝图不存在")
        base = self._blueprint_row_from_db(row)
        inspection = SmartReportAIInspectionResponse(**json.loads(str(row[3] or "{}")))
        return SmartReportBlueprintDetail(**base.model_dump(), inspection=inspection)

    async def preview_blueprint(self, blueprint_id: int) -> SmartReportBlueprintPreviewResponse:
        detail = await self.get_blueprint(blueprint_id)
        preview_text = self._render_blueprint_preview(detail.inspection)
        warnings = ["存在待确认项，建议确认后再生成正式报告。"] if detail.issue_count else []
        return SmartReportBlueprintPreviewResponse(
            blueprint_id=blueprint_id,
            preview_text=preview_text,
            issue_count=detail.issue_count,
            warnings=warnings,
        )

    async def generate_blueprint(self, blueprint_id: int) -> SmartReportBlueprintGenerateResponse:
        detail = await self.get_blueprint(blueprint_id)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        output_filename = f"smart_report_blueprint_{blueprint_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        output_path = self.output_dir / output_filename
        self._write_blueprint_docx(output_path, detail)
        finished = _iso_now()
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute(
                """
                UPDATE smart_report_blueprint
                SET output_file_path = ?, last_generated_at = ?, updated_at = ?, status = 'confirmed'
                WHERE blueprint_id = ?
                """,
                (str(output_path), finished, finished, blueprint_id),
            )
            await db.commit()
        return SmartReportBlueprintGenerateResponse(
            blueprint_id=blueprint_id,
            output_filename=output_filename,
            download_url=f"/api/smart-reports/blueprints/{blueprint_id}/download",
            generated_at=finished,
        )

    async def blueprint_output_path(self, blueprint_id: int) -> Path:
        detail = await self.get_blueprint(blueprint_id)
        if not detail.output_file_path:
            raise HTTPException(status_code=404, detail="报告蓝图尚未生成 Word")
        path = Path(detail.output_file_path)
        if not path.exists() or not path.is_file():
            raise HTTPException(status_code=404, detail="报告蓝图文件不存在")
        return path

    async def create_or_update_template(
        self,
        *,
        file: UploadFile,
        template_code: str,
        template_name: str,
        template_type: str = "analysis",
        remark: str | None = None,
        created_by: str | None = None,
    ) -> SmartReportTemplateCreateResponse:
        if not file.filename:
            raise HTTPException(status_code=400, detail="请上传 .docx 或 .pptx 文件")
        suffix = Path(file.filename).suffix.lower()
        if suffix not in {".docx", ".pptx"}:
            raise HTTPException(status_code=400, detail="请上传 .docx 或 .pptx 文件")
        safe_code = _safe_code(template_code)
        now = _iso_now()
        self.template_dir.mkdir(parents=True, exist_ok=True)

        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                "SELECT template_id, version_no FROM smart_report_template WHERE template_code = ?",
                (safe_code,),
            )
            existing = await cur.fetchone()
            next_version = int(existing[1]) + 1 if existing else 1
            target = self.template_dir / f"{safe_code}_v{next_version}{suffix}"
            with target.open("wb") as fh:
                shutil.copyfileobj(file.file, fh)

            placeholders = self._extract_placeholders_for_path(target)
            if existing:
                template_id = int(existing[0])
                await db.execute(
                    """
                    UPDATE smart_report_template
                    SET template_name = ?, template_type = ?, file_path = ?, version_no = ?,
                        remark = ?, updated_at = ?
                    WHERE template_id = ?
                    """,
                    (template_name.strip(), template_type.strip() or "analysis", str(target), next_version, remark, now, template_id),
                )
            else:
                cur = await db.execute(
                    """
                    INSERT INTO smart_report_template (
                      template_code, template_name, template_type, file_path, status,
                      version_no, remark, created_by, created_at, updated_at
                    ) VALUES (?, ?, ?, ?, 'active', ?, ?, ?, ?, ?)
                    """,
                    (
                        safe_code,
                        template_name.strip() or safe_code,
                        template_type.strip() or "analysis",
                        str(target),
                        next_version,
                        remark,
                        created_by,
                        now,
                        now,
                    ),
                )
                template_id = int(cur.lastrowid)

            await self._sync_detected_variables(db, template_id, placeholders, now)
            await db.commit()

        template = await self.get_template(template_id)
        return SmartReportTemplateCreateResponse(template=template, placeholders=placeholders)

    async def create_or_update_text_template(
        self,
        *,
        template_code: str,
        template_name: str,
        content: str,
        template_type: str = "analysis",
        remark: str | None = None,
        created_by: str | None = None,
    ) -> SmartReportTemplateCreateResponse:
        safe_code = _safe_code(template_code)
        now = _iso_now()
        self.template_dir.mkdir(parents=True, exist_ok=True)

        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                "SELECT template_id, version_no FROM smart_report_template WHERE template_code = ?",
                (safe_code,),
            )
            existing = await cur.fetchone()
            next_version = int(existing[1]) + 1 if existing else 1
            target = self.template_dir / f"{safe_code}_v{next_version}.docx"
            self._write_text_template_docx(target, template_name.strip() or safe_code, content)
            placeholders = self.extract_placeholders(target)

            if existing:
                template_id = int(existing[0])
                await db.execute(
                    """
                    UPDATE smart_report_template
                    SET template_name = ?, template_type = ?, file_path = ?, version_no = ?,
                        remark = ?, updated_at = ?
                    WHERE template_id = ?
                    """,
                    (template_name.strip(), template_type.strip() or "analysis", str(target), next_version, remark, now, template_id),
                )
            else:
                cur = await db.execute(
                    """
                    INSERT INTO smart_report_template (
                      template_code, template_name, template_type, file_path, status,
                      version_no, remark, created_by, created_at, updated_at
                    ) VALUES (?, ?, ?, ?, 'active', ?, ?, ?, ?, ?)
                    """,
                    (
                        safe_code,
                        template_name.strip() or safe_code,
                        template_type.strip() or "analysis",
                        str(target),
                        next_version,
                        remark,
                        created_by,
                        now,
                        now,
                    ),
                )
                template_id = int(cur.lastrowid)

            await self._sync_detected_variables(db, template_id, placeholders, now)
            await db.commit()

        template = await self.get_template(template_id)
        return SmartReportTemplateCreateResponse(template=template, placeholders=placeholders)

    async def get_template(self, template_id: int) -> SmartReportTemplateRow:
        rows = [t for t in await self.list_templates() if t.template_id == template_id]
        if not rows:
            raise HTTPException(status_code=404, detail="报告模板不存在")
        return rows[0]

    async def list_variables(self, template_id: int) -> list[SmartReportTemplateVariableRow]:
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                """
                SELECT variable_id, template_id, variable_key, variable_name, variable_type,
                       binding_config_json, display_order, created_at, updated_at
                FROM smart_report_template_variable
                WHERE template_id = ?
                ORDER BY display_order, variable_id
                """,
                (template_id,),
            )
            rows = await cur.fetchall()
        return [
            SmartReportTemplateVariableRow(
                variable_id=int(r[0]),
                template_id=int(r[1]),
                variable_key=str(r[2]),
                variable_name=str(r[3]),
                variable_type=str(r[4]),  # type: ignore[arg-type]
                binding_config=_json_loads_dict(r[5]),
                display_order=int(r[6] or 0),
                created_at=str(r[7]),
                updated_at=str(r[8]),
            )
            for r in rows
        ]

    async def upsert_variables(
        self, template_id: int, variables: list[SmartReportTemplateVariableUpsert]
    ) -> list[SmartReportTemplateVariableRow]:
        await self.get_template(template_id)
        now = _iso_now()
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            confirmed_codes = await load_confirmed_org_product_runtime_ref_codes(db)
            for item in variables:
                inferred_type, inferred_config = self._infer_variable(item.variable_key)
                variable_type = item.variable_type or inferred_type
                binding_config = item.binding_config or inferred_config
                self._validate_formula_binding_config(
                    variable_type,
                    binding_config,
                    confirmed_codes=confirmed_codes,
                )
                variable_name = item.variable_name or self._default_variable_name(item.variable_key)
                await db.execute(
                    """
                    INSERT INTO smart_report_template_variable (
                      template_id, variable_key, variable_name, variable_type,
                      binding_config_json, display_order, created_at, updated_at
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                    ON CONFLICT(template_id, variable_key) DO UPDATE SET
                      variable_name = excluded.variable_name,
                      variable_type = excluded.variable_type,
                      binding_config_json = excluded.binding_config_json,
                      display_order = excluded.display_order,
                      updated_at = excluded.updated_at
                    """,
                    (
                        template_id,
                        item.variable_key.strip(),
                        variable_name,
                        variable_type,
                        json.dumps(binding_config, ensure_ascii=False),
                        item.display_order,
                        now,
                        now,
                    ),
                )
            await db.commit()
        return await self.list_variables(template_id)

    async def list_calc_metrics(self) -> list[SmartReportCalcMetricRow]:
        async with aiosqlite.connect(common_db_path()) as db:
            cur = await db.execute(
                """
                SELECT metric_code, metric_name, expression, components_json, value_type,
                       format_type, remark, created_at, updated_at
                FROM smart_report_calc_metric
                ORDER BY updated_at DESC, metric_code
                """
            )
            rows = await cur.fetchall()
        return [self._calc_metric_from_row(row) for row in rows]

    async def upsert_calc_metric(self, body: SmartReportCalcMetricUpsert) -> SmartReportCalcMetricRow:
        metric_code = re.sub(r"[^A-Za-z0-9_]+", "_", body.metric_code.strip()).strip("_")
        if not metric_code:
            raise HTTPException(status_code=400, detail="计算指标编码不能为空")
        aliases = {item.alias.strip() for item in body.components}
        if not aliases:
            raise HTTPException(status_code=400, detail="请至少选择一个基础指标")
        self._safe_eval_expression(body.expression, {alias: 1.0 for alias in aliases})
        now = _iso_now()
        components = [item.model_dump() for item in body.components]
        async with aiosqlite.connect(common_db_path()) as db:
            confirmed_codes = await load_confirmed_org_product_runtime_ref_codes(db)
            component_codes = {
                str(item.get("data_acct_code") or "").strip().upper()
                for item in components
                if str(item.get("data_acct_code") or "").strip()
            }
            unconfirmed_codes = sorted(component_codes - confirmed_codes)
            if unconfirmed_codes:
                raise HTTPException(
                    status_code=400,
                    detail=(
                        "计算指标组成项未在机构及产品指标主表中确认："
                        f"{'、'.join(unconfirmed_codes[:10])}"
                    ),
                )
            await db.execute(
                """
                INSERT INTO smart_report_calc_metric (
                  metric_code, metric_name, expression, components_json, value_type,
                  format_type, remark, created_at, updated_at
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(metric_code) DO UPDATE SET
                  metric_name = excluded.metric_name,
                  expression = excluded.expression,
                  components_json = excluded.components_json,
                  value_type = excluded.value_type,
                  format_type = excluded.format_type,
                  remark = excluded.remark,
                  updated_at = excluded.updated_at
                """,
                (
                    metric_code,
                    body.metric_name.strip(),
                    body.expression.strip(),
                    json.dumps(components, ensure_ascii=False),
                    body.value_type.strip() or "金额",
                    body.format_type.strip() or "number",
                    body.remark,
                    now,
                    now,
                ),
            )
            await db.commit()
        metric = await self.get_calc_metric(metric_code)
        if not metric:
            raise HTTPException(status_code=500, detail="计算指标保存失败")
        return metric

    async def get_calc_metric(self, metric_code: str) -> SmartReportCalcMetricRow | None:
        async with aiosqlite.connect(common_db_path()) as db:
            cur = await db.execute(
                """
                SELECT metric_code, metric_name, expression, components_json, value_type,
                       format_type, remark, created_at, updated_at
                FROM smart_report_calc_metric
                WHERE metric_code = ?
                """,
                (metric_code.strip(),),
            )
            row = await cur.fetchone()
        return self._calc_metric_from_row(row) if row else None

# ─── 报告生成与渲染 ───

    async def generate(self, body: SmartReportGenerateRequest) -> SmartReportGenerateResponse:
        now = _iso_now()
        template_row = await self._fetch_template_record(body.template_id)
        instance_name = body.instance_name or f"{template_row['template_name']} {now}"
        self.output_dir.mkdir(parents=True, exist_ok=True)

        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                """
                INSERT INTO smart_report_instance (
                  template_id, instance_name, parameter_values_json,
                  text_values_json, generation_status, created_at, updated_at
                ) VALUES (?, ?, ?, ?, 'running', ?, ?)
                """,
                (
                    body.template_id,
                    instance_name,
                    json.dumps(body.parameters, ensure_ascii=False),
                    json.dumps(body.text_values, ensure_ascii=False),
                    now,
                    now,
                ),
            )
            instance_id = int(cur.lastrowid)
            cur = await db.execute(
                """
                INSERT INTO smart_report_job (instance_id, job_type, job_status, started_at)
                VALUES (?, 'generate', 'running', ?)
                """,
                (instance_id, now),
            )
            job_id = int(cur.lastrowid)
            await db.commit()

        return await self._render_instance(
            instance_id=instance_id,
            job_id=job_id,
            template_id=body.template_id,
            parameters=body.parameters,
            text_values=body.text_values,
            job_type="generate",
        )

    async def preview(self, body: SmartReportPreviewRequest) -> SmartReportPreviewResponse:
        template_row = await self._fetch_template_record(body.template_id)
        template_path = Path(str(template_row["file_path"]))
        variables = await self._variables_for_template(body.template_id, template_path)
        resolved, warnings = await self._resolve_values(variables, body.parameters, body.text_values)
        chart_labels = await self._chart_preview_labels(template_path)
        return SmartReportPreviewResponse(
            preview_text=self._render_preview_text(template_path, resolved, chart_labels=chart_labels),
            resolved_values=resolved,
            warnings=warnings,
        )

    async def refresh_instance(self, instance_id: int) -> SmartReportGenerateResponse:
        now = _iso_now()
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                """
                SELECT template_id, parameter_values_json, text_values_json
                FROM smart_report_instance
                WHERE instance_id = ?
                """,
                (instance_id,),
            )
            row = await cur.fetchone()
            if not row:
                raise HTTPException(status_code=404, detail="报告实例不存在")
            template_id = int(row[0])
            parameters = _json_loads_dict(row[1])
            text_values = _json_loads_dict(row[2])
            await db.execute(
                """
                UPDATE smart_report_instance
                SET generation_status = 'running', error_message = NULL, updated_at = ?
                WHERE instance_id = ?
                """,
                (now, instance_id),
            )
            cur = await db.execute(
                """
                INSERT INTO smart_report_job (instance_id, job_type, job_status, started_at)
                VALUES (?, 'refresh', 'running', ?)
                """,
                (instance_id, now),
            )
            job_id = int(cur.lastrowid)
            await db.commit()

        return await self._render_instance(
            instance_id=instance_id,
            job_id=job_id,
            template_id=template_id,
            parameters=parameters,
            text_values=text_values,
            job_type="refresh",
        )

    async def _render_instance(
        self,
        *,
        instance_id: int,
        job_id: int,
        template_id: int,
        parameters: dict[str, Any],
        text_values: dict[str, Any],
        job_type: str,
    ) -> SmartReportGenerateResponse:
        try:
            self.output_dir.mkdir(parents=True, exist_ok=True)
            template_row = await self._fetch_template_record(template_id)
            template_path = Path(str(template_row["file_path"]))
            variables = await self._variables_for_template(template_id, template_path)
            resolved, warnings = await self._resolve_values(variables, parameters, text_values)
            prefix = "smart_report_refresh" if job_type == "refresh" else "smart_report"
            ext = template_path.suffix.lower()
            output_filename = f"{prefix}_{instance_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{ext}"
            output_path = self.output_dir / output_filename
            if ext == ".pptx":
                self._render_pptx(template_path, output_path, resolved)
            else:
                await self._render_docx(template_path, output_path, resolved, parameters)
            finished = _iso_now()
            async with aiosqlite.connect(common_db_path()) as db:
                await db.execute("PRAGMA foreign_keys = ON")
                if job_type == "refresh":
                    await db.execute(
                        """
                        UPDATE smart_report_instance
                        SET output_file_path = ?, data_snapshot_json = ?, generation_status = 'success',
                            last_refresh_at = ?, updated_at = ?
                        WHERE instance_id = ?
                        """,
                        (
                            str(output_path),
                            json.dumps(resolved, ensure_ascii=False),
                            finished,
                            finished,
                            instance_id,
                        ),
                    )
                else:
                    await db.execute(
                        """
                        UPDATE smart_report_instance
                        SET output_file_path = ?, data_snapshot_json = ?, generation_status = 'success',
                            last_generated_at = ?, updated_at = ?
                        WHERE instance_id = ?
                        """,
                        (
                            str(output_path),
                            json.dumps(resolved, ensure_ascii=False),
                            finished,
                            finished,
                            instance_id,
                        ),
                    )
                await db.execute(
                    """
                    UPDATE smart_report_job
                    SET job_status = 'success', finished_at = ?
                    WHERE job_id = ?
                    """,
                    (finished, job_id),
                )
                await db.commit()
            return SmartReportGenerateResponse(
                instance_id=instance_id,
                job_id=job_id,
                output_filename=output_filename,
                download_url=f"/api/smart-reports/instances/{instance_id}/download",
                generated_at=finished,
                resolved_values=resolved,
                warnings=warnings,
            )
        except Exception as exc:
            finished = _iso_now()
            async with aiosqlite.connect(common_db_path()) as db:
                await db.execute("PRAGMA foreign_keys = ON")
                await db.execute(
                    """
                    UPDATE smart_report_instance
                    SET generation_status = 'failed', error_message = ?, updated_at = ?
                    WHERE instance_id = ?
                    """,
                    (str(exc), finished, instance_id),
                )
                await db.execute(
                    """
                    UPDATE smart_report_job
                    SET job_status = 'failed', finished_at = ?, error_message = ?
                    WHERE job_id = ?
                    """,
                    (finished, str(exc), job_id),
                )
                await db.commit()
            if isinstance(exc, HTTPException):
                raise exc
            raise HTTPException(status_code=500, detail=f"生成报告失败：{exc}") from exc

    async def list_instances(self) -> list[SmartReportInstanceRow]:
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                """
                SELECT i.instance_id, i.template_id, t.template_name,
                       i.instance_name, i.generation_status, i.output_file_path,
                       i.error_message, i.last_generated_at, i.last_refresh_at,
                       i.created_at, i.updated_at
                FROM smart_report_instance i
                LEFT JOIN smart_report_template t ON t.template_id = i.template_id
                ORDER BY i.updated_at DESC, i.instance_id DESC
                LIMIT 100
                """
            )
            rows = await cur.fetchall()
        return [
            SmartReportInstanceRow(
                instance_id=int(r[0]),
                template_id=int(r[1]),
                template_name=str(r[2]) if r[2] is not None else None,
                instance_name=str(r[3]),
                generation_status=str(r[4]),
                output_file_path=str(r[5]) if r[5] is not None else None,
                error_message=str(r[6]) if r[6] is not None else None,
                last_generated_at=str(r[7]) if r[7] is not None else None,
                last_refresh_at=str(r[8]) if r[8] is not None else None,
                created_at=str(r[9]),
                updated_at=str(r[10]),
            )
            for r in rows
        ]

    async def instance_output_path(self, instance_id: int) -> Path:
        async with aiosqlite.connect(common_db_path()) as db:
            cur = await db.execute(
                "SELECT output_file_path FROM smart_report_instance WHERE instance_id = ?",
                (instance_id,),
            )
            row = await cur.fetchone()
        if not row or not row[0]:
            raise HTTPException(status_code=404, detail="报告实例没有可下载文件")
        path = Path(str(row[0]))
        if not path.exists() or not path.is_file():
            raise HTTPException(status_code=404, detail="报告文件不存在")
        return path

    def extract_placeholders(self, docx_path: Path) -> list[str]:
        doc = Document(str(docx_path))
        found: list[str] = []
        for text in self._iter_doc_text(doc):
            for match in PLACEHOLDER_RE.finditer(text):
                token = match.group(1).strip()
                if token and token not in found:
                    found.append(token)
        return found

    def extract_placeholders_pptx(self, pptx_path: Path) -> list[str]:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        found: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for match in PLACEHOLDER_RE.finditer(paragraph.text):
                            token = match.group(1).strip()
                            if token and token not in found:
                                found.append(token)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for match in PLACEHOLDER_RE.finditer(cell.text):
                                token = match.group(1).strip()
                                if token and token not in found:
                                    found.append(token)
        return found

    async def _sync_detected_variables(
        self, db: aiosqlite.Connection, template_id: int, placeholders: list[str], now: str
    ) -> None:
        confirmed_codes = await load_confirmed_org_product_runtime_ref_codes(db)
        for idx, token in enumerate(placeholders):
            variable_type, binding_config = self._infer_variable(token)
            self._validate_formula_binding_config(
                variable_type,
                binding_config,
                confirmed_codes=confirmed_codes,
            )
            await db.execute(
                """
                INSERT INTO smart_report_template_variable (
                  template_id, variable_key, variable_name, variable_type,
                  binding_config_json, display_order, created_at, updated_at
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                ON CONFLICT(template_id, variable_key) DO UPDATE SET
                  display_order = excluded.display_order,
                  updated_at = excluded.updated_at
                """,
                (
                    template_id,
                    token,
                    self._default_variable_name(token),
                    variable_type,
                    json.dumps(binding_config, ensure_ascii=False),
                    idx,
                    now,
                    now,
                ),
            )

    def _validate_formula_binding_config(
        self,
        variable_type: str,
        binding_config: dict[str, Any],
        *,
        confirmed_codes: set[str],
    ) -> None:
        if str(variable_type or "").strip().lower() != "formula":
            return
        data_acct_code = str(binding_config.get("data_acct_code") or "").strip().upper()
        if not data_acct_code:
            return
        if data_acct_code not in confirmed_codes:
            raise HTTPException(
                status_code=400,
                detail=f"报告公式变量未在机构及产品指标主表中确认：{data_acct_code}",
            )

    def _infer_variable(self, token: str) -> tuple[str, dict[str, Any]]:
        text = token.strip()
        if ":" in text:
            prefix, raw_key = text.split(":", 1)
            prefix = prefix.strip().lower()
            key = raw_key.strip()
            if prefix == "metric":
                return "metric", {"metric_id": key}
            if prefix == "formula":
                parts = [p.strip() for p in key.split(":") if p.strip()]
                data_acct_code = parts[0] if parts else key
                formula_type = parts[1] if len(parts) > 1 else "budget"
                return "formula", {"data_acct_code": data_acct_code, "formula_type": formula_type}
            if prefix in {"calc", "calculated", "calculated_metric"}:
                return "calc", {"metric_code": key}
            if prefix in {"param", "parameter"}:
                return "parameter", {"param_key": key}
            if prefix in {"table", "chart", "analysis"}:
                return prefix, {f"{prefix}_id": key}
            if prefix == "text":
                return "text", {"text_key": key}
        return "text", {"text_key": text}

    def _default_variable_name(self, token: str) -> str:
        if ":" in token:
            prefix, key = token.split(":", 1)
            prefix_name = {
                "metric": "指标",
                "formula": "公式",
                "calc": "计算指标",
                "calculated": "计算指标",
                "calculated_metric": "计算指标",
                "param": "参数",
                "parameter": "参数",
                "table": "表格",
                "chart": "图表",
                "analysis": "分析结论",
                "text": "文本",
            }.get(prefix.strip().lower(), "变量")
            return f"{prefix_name} {key.strip()}"
        return token.strip()

    async def _fetch_template_record(self, template_id: int) -> dict[str, Any]:
        async with aiosqlite.connect(common_db_path()) as db:
            await db.execute("PRAGMA foreign_keys = ON")
            cur = await db.execute(
                "SELECT template_id, template_name, file_path FROM smart_report_template WHERE template_id = ?",
                (template_id,),
            )
            row = await cur.fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="报告模板不存在")
        return {"template_id": int(row[0]), "template_name": str(row[1]), "file_path": str(row[2])}

    async def _resolve_values(
        self,
        variables: list[SmartReportTemplateVariableRow],
        parameters: dict[str, Any],
        text_values: dict[str, Any],
    ) -> tuple[dict[str, str], list[str]]:
        resolved: dict[str, str] = {}
        warnings: list[str] = []
        for variable in variables:
            token = variable.variable_key
            inferred_type, inferred_config = self._infer_variable(token)
            explicit_prefix = token.split(":", 1)[0].strip().lower() if ":" in token else ""
            variable_type = inferred_type if explicit_prefix in {"metric", "formula", "calc", "calculated", "calculated_metric", "param", "parameter", "text", "table", "chart", "analysis"} else variable.variable_type
            cfg = {**inferred_config, **(variable.binding_config or {})}
            if variable_type == "parameter":
                key = str(cfg.get("param_key") or token.split(":", 1)[-1]).strip()
                resolved[token] = str(parameters.get(key, ""))
            elif variable_type == "metric":
                metric_id = str(cfg.get("metric_id") or token.split(":", 1)[-1]).strip()
                value = await self._calculate_metric(metric_id, parameters)
                resolved[token] = value
            elif variable_type == "formula":
                resolved[token] = await self._resolve_formula_text(token, cfg, parameters)
            elif variable_type == "calc":
                metric_code = str(cfg.get("metric_code") or token.split(":", 1)[-1]).strip()
                resolved[token] = await self._calculate_report_metric(metric_code, parameters)
            elif variable_type == "text":
                key = str(cfg.get("text_key") or token.split(":", 1)[-1]).strip()
                resolved[token] = str(text_values.get(key, parameters.get(key, "")))
            elif variable_type == "chart":
                resolved[token] = ""
            else:
                resolved[token] = f"[{variable.variable_name} 待接入]"
                warnings.append(f"{variable.variable_key} 当前为占位输出，后续接入 {variable_type} 渲染")
        return resolved, warnings

    async def _resolve_formula_text(self, token: str, cfg: dict[str, Any], params: dict[str, Any]) -> str:
        raw = token.split(":", 1)[-1] if ":" in token else token
        parts = [p.strip() for p in raw.split(":") if p.strip()]
        data_acct_code = str(cfg.get("data_acct_code") or (parts[0] if parts else "")).strip().upper()
        formula_type = str(cfg.get("formula_type") or (parts[1] if len(parts) > 1 else "auto")).strip().lower()
        if formula_type == "auto":
            formula_type = "actual" if self._budget_actual_param(params) == 1 else "budget"
        if not data_acct_code:
            return ""
        async with aiosqlite.connect(common_db_path()) as db:
            confirmed_codes = await load_confirmed_org_product_runtime_ref_codes(db)
            self._validate_formula_binding_config(
                "formula",
                {"data_acct_code": data_acct_code},
                confirmed_codes=confirmed_codes,
            )
            cur = await db.execute(
                """
                SELECT data_acct_code, data_acct_name, budget_formula, actual_formula
                FROM data_account
                WHERE data_acct_code = ?
                """,
                (data_acct_code,),
            )
            row = await cur.fetchone()
        if not row:
            return f"{data_acct_code} 未找到公式"
        formula = row[3] if formula_type in {"actual", "1", "实际"} else row[2]
        label = "实际公式" if formula_type in {"actual", "1", "实际"} else "预算公式"
        return f"{row[0]} {row[1]} {label}：{formula or '未配置'}"

    async def _calculate_report_metric(self, metric_code: str, params: dict[str, Any]) -> str:
        metric = await self.get_calc_metric(metric_code)
        if not metric:
            return f"{metric_code} 未找到计算指标"
        values: dict[str, float] = {}
        for component in metric.components:
            values[component.alias] = await self._sum_data_account(component.data_acct_code, params, self._budget_actual_param(params))
        try:
            result = self._safe_eval_expression(metric.expression, values)
        except ZeroDivisionError:
            return "N/A"
        percentage = metric.format_type == "percent" or metric.value_type == "百分比"
        return _format_number(result, percentage=percentage)

    async def _calculate_metric(self, metric_id: str, params: dict[str, Any]) -> str:
        if metric_id == "m_budget_actual_gap":
            actual = await self._sum_budget_summary(params, budget_actual=1)
            budget = await self._sum_budget_summary(params, budget_actual=0)
            return _format_number(actual - budget)
        if metric_id == "m_yoy_growth":
            current = await self._sum_budget_summary(params, budget_actual=self._budget_actual_param(params))
            previous_params = dict(params)
            previous_params["year"] = _plain_year(params.get("year")) - 1
            previous = await self._sum_budget_summary(previous_params, budget_actual=self._budget_actual_param(params))
            if previous == 0:
                return "N/A"
            return _format_number((current - previous) / previous, percentage=True)
        if metric_id == "m_mom_growth":
            current_month = _month_label(params.get("month"))
            if not current_month:
                return "N/A"
            month_num = int(current_month[1:])
            if month_num <= 1:
                return "N/A"
            current = await self._sum_budget_summary(params, budget_actual=self._budget_actual_param(params))
            previous_params = dict(params)
            previous_params["month"] = f"M{month_num - 1:02d}"
            previous = await self._sum_budget_summary(previous_params, budget_actual=self._budget_actual_param(params))
            if previous == 0:
                return "N/A"
            return _format_number((current - previous) / previous, percentage=True)
        return "N/A"

    async def _sum_data_account(self, data_acct_code: str, params: dict[str, Any], budget_actual: int) -> float:
        scoped_params = dict(params)
        scoped_params["data_acct_code"] = data_acct_code
        return await self._sum_budget_summary(scoped_params, budget_actual=budget_actual)

    async def _load_ai_binding_catalog(self, limit: int = 180) -> list[dict[str, Any]]:
        async with aiosqlite.connect(common_db_path()) as db:
            confirmed_codes = sorted(await load_confirmed_org_product_runtime_ref_codes(db))
            if not confirmed_codes:
                return []
            placeholders = ",".join("?" for _ in confirmed_codes)
            cur = await db.execute(
                f"""
                SELECT d.data_acct_code, d.data_acct_name,
                       d.value_type, d.budget_formula, d.actual_formula,
                       GROUP_CONCAT(n.node_name, ' / ') AS metric_nodes
                FROM data_account d
                LEFT JOIN data_account_metric_binding b ON b.data_acct_code = d.data_acct_code
                LEFT JOIN data_account_metric_node n ON n.node_code = b.metric_node_code
                WHERE UPPER(d.data_acct_code) IN ({placeholders})
                GROUP BY d.data_acct_code
                ORDER BY d.data_acct_code
                LIMIT ?
                """,
                (*confirmed_codes, limit),
            )
            rows = await cur.fetchall()
        return [
            {
                "code": str(row[0]),
                "name": str(row[1] or ""),
                "group": str(row[5] or ""),
                "value_type": str(row[2] or ""),
                "budget_formula": str(row[3] or "")[:120],
                "actual_formula": str(row[4] or "")[:120],
                "metric_nodes": str(row[5] or "")[:160],
            }
            for row in rows
        ]

    def _call_ai_report_inspector(self, report_text: str, catalog: list[dict[str, Any]]) -> str | None:
        if not self.deepseek_client or not self.deepseek_client.is_enabled():
            return None
        system_prompt = "你是银行预算经营分析报告结构化专家。请把上传的 Word 报告理解成可刷新报告蓝图。只输出严格 JSON，不要 Markdown。"
        user_prompt = json.dumps(
            {
                "task": "解析报告结构，识别指标、公式说明、自然语言分析规则，并给出需要用户确认的问题项。",
                "output_schema": {
                    "summary": "一句话概括报告主题",
                    "blocks": [
                        {
                            "block_id": "B1",
                            "block_type": "text_block | metric_value | formula_explain | analysis_task | table_block | unmatched_item",
                            "text": "原文片段",
                            "metrics": [
                                {
                                    "name": "报告中的指标名",
                                    "matched_code": "匹配到的机构及产品指标编码，未匹配为空",
                                    "matched_name": "匹配名称",
                                    "confidence": 0.0,
                                    "reason": "匹配依据",
                                }
                            ],
                            "analysis_rule_nl": "如果是同比归因/TopN/预算偏差等分析任务，保留自然语言规则",
                            "structured_plan": {
                                "compare_type": "yoy | mom | budget_vs_actual | unknown",
                                "metric": "分析指标",
                                "dimension": "product | dept | data_account | unknown",
                                "direction": "increase | decrease | variance | unknown",
                                "top_n": 3,
                                "output": "summary_text | table | summary_text_and_table",
                            },
                            "confidence": 0.0,
                        }
                    ],
                    "issues": [
                        {
                            "issue_type": "unmatched_metric | low_confidence_match | analysis_rule_needs_confirmation | caliber_conflict",
                            "text": "需要用户确认的问题",
                            "suggested_action": "建议用户怎么补充",
                            "candidates": [
                                {
                                    "target_type": "data_account | metric_node | analysis_plan",
                                    "target_code": "候选编码",
                                    "target_name": "候选名称",
                                    "confidence": 0.0,
                                    "reason": "推荐原因",
                                }
                            ],
                            "rule_preview": "把自然语言规则翻译成业务可读执行计划",
                        }
                    ],
                    "assumptions": ["模型做出的假设"],
                },
                "matching_catalog": catalog,
                "report_text": report_text,
                "rules": [
                    "高置信度匹配不要制造问题项；低置信度或无匹配必须进入 issues。",
                    "同比归因、规模下降前三产品、预算偏差原因等不要当成普通指标，要识别为 analysis_task。",
                    "自然语言分析规则必须给 rule_preview，便于用户确认。",
                ],
            },
            ensure_ascii=False,
        )
        return self.deepseek_client.chat_completion(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            temperature=0.1,
            max_tokens=2200,
        )

    def _parse_ai_inspection(self, raw: str) -> dict[str, Any]:
        text = raw.strip()
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
        try:
            parsed = json.loads(text)
        except Exception as exc:
            raise HTTPException(status_code=502, detail=f"AI 返回结果不是合法 JSON：{exc}") from exc
        if not isinstance(parsed, dict):
            raise HTTPException(status_code=502, detail="AI 返回结果结构错误")
        return parsed

    def _fallback_ai_inspection(self, compact_text: str) -> dict[str, Any]:
        sentences = [s.strip() for s in re.split(r"[。；;\n]", compact_text) if s.strip()]
        blocks: list[dict[str, Any]] = []
        issues: list[dict[str, Any]] = []
        for idx, sentence in enumerate(sentences[:20]):
            is_analysis = any(key in sentence for key in ["同比", "环比", "归因", "前三", "前3", "下降", "上升", "偏差"])
            block: dict[str, Any] = {
                "block_id": f"B{idx + 1}",
                "block_type": "analysis_task" if is_analysis else "text_block",
                "text": sentence[:300],
                "metrics": [],
                "confidence": 0.35 if is_analysis else 0.2,
            }
            if is_analysis:
                block["analysis_rule_nl"] = sentence
                block["structured_plan"] = {
                    "compare_type": "yoy" if "同比" in sentence else "unknown",
                    "metric": "",
                    "dimension": "product" if "产品" in sentence else "unknown",
                    "direction": "decrease" if "下降" in sentence else "increase" if "上升" in sentence else "variance",
                    "top_n": 3 if ("前三" in sentence or "前3" in sentence) else None,
                    "output": "summary_text_and_table",
                }
                issues.append(
                    {
                        "issue_type": "analysis_rule_needs_confirmation",
                        "text": sentence[:220],
                        "suggested_action": "请确认分析指标、比较口径、维度和 TopN 数量。",
                        "candidates": [],
                        "rule_preview": "系统将根据你补充的自然语言规则生成可执行分析计划。",
                    }
                )
            blocks.append(block)
        return {
            "summary": sentences[0][:120] if sentences else "报告待解析",
            "blocks": blocks,
            "issues": issues,
            "assumptions": ["规则兜底解析只识别明显的同比/环比/归因/TopN 语句，建议启用 DeepSeek 获得更准确结果。"],
        }

    def _extract_docx_plain_text(self, docx_path: Path) -> str:
        doc = Document(str(docx_path))
        lines: list[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)

    def _blueprint_row_from_db(self, row: Any) -> SmartReportBlueprintRow:
        inspection = _json_loads_dict(str(row[3] or "{}"))
        blocks = inspection.get("blocks") if isinstance(inspection.get("blocks"), list) else []
        issues = inspection.get("issues") if isinstance(inspection.get("issues"), list) else []
        return SmartReportBlueprintRow(
            blueprint_id=int(row[0]),
            blueprint_name=str(row[1]),
            source_filename=str(row[2]),
            status=str(row[4] or "draft"),
            issue_count=len(issues),
            block_count=len(blocks),
            output_file_path=str(row[5]) if row[5] else None,
            last_generated_at=str(row[6]) if row[6] else None,
            created_at=str(row[7]),
            updated_at=str(row[8]),
        )

    def _render_blueprint_preview(self, inspection: SmartReportAIInspectionResponse) -> str:
        lines: list[str] = []
        lines.append(f"报告摘要：{inspection.summary or '未生成摘要'}")
        lines.append("")
        lines.append("一、AI 识别的报告结构")
        for block in inspection.blocks:
            lines.append(f"[{block.block_id}] {block.block_type} 置信度 {block.confidence:.2f}")
            if block.text:
                lines.append(block.text)
            if block.metrics:
                for metric in block.metrics:
                    name = metric.get("name") or ""
                    code = metric.get("matched_code") or "未匹配"
                    matched_name = metric.get("matched_name") or ""
                    confidence = metric.get("confidence", "")
                    lines.append(f"  - 指标：{name} -> {code} {matched_name} 置信度 {confidence}")
            if block.analysis_rule_nl:
                lines.append(f"  - 自然语言规则：{block.analysis_rule_nl}")
            if block.structured_plan:
                lines.append(f"  - 执行计划：{self._business_plan_text(block.structured_plan)}")
            lines.append("")
        lines.append("二、待确认项")
        if inspection.issues:
            for idx, issue in enumerate(inspection.issues, start=1):
                lines.append(f"{idx}. {issue.issue_type}：{issue.text}")
                if issue.rule_preview:
                    lines.append(f"   计划预览：{issue.rule_preview}")
                if issue.suggested_action:
                    lines.append(f"   建议动作：{issue.suggested_action}")
        else:
            lines.append("无待确认项。")
        if inspection.assumptions:
            lines.append("")
            lines.append("三、AI 假设")
            lines.extend(f"- {item}" for item in inspection.assumptions)
        return "\n".join(lines)

    def _business_plan_text(self, plan: dict[str, Any]) -> str:
        compare = {"yoy": "同比", "mom": "环比", "budget_vs_actual": "预算实际偏差"}.get(str(plan.get("compare_type") or ""), "待确认比较口径")
        dimension = {"product": "产品", "dept": "部门", "data_account": "机构及产品指标编码"}.get(
            str(plan.get("dimension") or ""),
            "待确认维度",
        )
        direction = {"increase": "上升", "decrease": "下降", "variance": "偏差"}.get(str(plan.get("direction") or ""), "波动")
        top_n = plan.get("top_n") or 3
        metric = plan.get("metric") or "待确认指标"
        return f"按{dimension}维度，对{metric}做{compare}{direction}归因，输出 Top {top_n} 及说明。"

    def _write_blueprint_docx(self, output_path: Path, detail: SmartReportBlueprintDetail) -> None:
        doc = Document()
        doc.add_heading(detail.blueprint_name, level=1)
        doc.add_paragraph(f"来源报告：{detail.source_filename}")
        doc.add_paragraph(f"生成时间：{_iso_now()}")
        doc.add_heading("AI 解析摘要", level=2)
        doc.add_paragraph(detail.inspection.summary or "未生成摘要")
        doc.add_heading("报告结构预览", level=2)
        for line in self._render_blueprint_preview(detail.inspection).splitlines():
            if not line.strip():
                doc.add_paragraph("")
            elif line.startswith("["):
                doc.add_paragraph(line, style="List Bullet")
            else:
                doc.add_paragraph(line)
        doc.save(str(output_path))

    def _budget_actual_param(self, params: dict[str, Any]) -> int:
        raw = str(params.get("budget_actual", params.get("caliber", "1"))).strip().lower()
        if raw in {"0", "budget", "预算"}:
            return 0
        return 1

    async def _sum_budget_summary(self, params: dict[str, Any], *, budget_actual: int) -> float:
        year_int = _plain_year(params.get("year"))
        path = budget_db_path(year_int)
        if not path.exists():
            return 0.0
        where = ["year = ?", "budget_actual = ?"]
        values: list[Any] = [_year_label(year_int), budget_actual]
        month = _month_label(params.get("month"))
        start_month = _month_label(params.get("start_month"))
        end_month = _month_label(params.get("end_month"))
        if start_month or end_month:
            start_num = int((start_month or "M01")[1:])
            end_num = int((end_month or "M12")[1:])
            if start_num > end_num:
                start_num, end_num = end_num, start_num
            where.append("CAST(SUBSTR(month, 2) AS INTEGER) BETWEEN ? AND ?")
            values.extend([start_num, end_num])
        elif month:
            where.append("month = ?")
            values.append(month)
        quarter = str(params.get("quarter") or "").strip().upper()
        if quarter:
            where.append("quarter = ?")
            values.append(quarter)
        version_id = params.get("version_id")
        if version_id is not None and str(version_id).strip() != "":
            where.append("version_id = ?")
            values.append(int(version_id))
        data_code = params.get("data_acct_code") or params.get("data_account")
        if data_code:
            where.append("data_code_name LIKE ?")
            values.append(f"%{str(data_code).strip()}%")
        product_code = params.get("product_code") or params.get("product")
        if product_code:
            where.append("IFNULL(product_code_name, '') LIKE ?")
            values.append(f"%{str(product_code).strip()}%")
        dept = params.get("dept_code") or params.get("dept")
        if dept:
            where.append("(IFNULL(dept_level1, '') || IFNULL(dept_level2, '') || IFNULL(dept_level3, '')) LIKE ?")
            values.append(f"%{str(dept).strip()}%")
        metric_code = params.get("metric_code") or params.get("metric")
        if metric_code:
            where.append(
                "(IFNULL(metric_level1, '') || IFNULL(metric_level2, '') || IFNULL(metric_level3, '') || IFNULL(metric_level4, '') || IFNULL(metric_level5, '')) LIKE ?"
            )
            values.append(f"%{str(metric_code).strip()}%")

        async with aiosqlite.connect(path) as db:
            await ensure_budget_summary_read_model_schema_async(db)
            where.append("value_source <> 'rollup'")
            cur = await db.execute(
                f"SELECT COALESCE(SUM(value), 0) FROM budget_summary WHERE {' AND '.join(where)}",
                values,
            )
            row = await cur.fetchone()
        return float(row[0] or 0.0) if row else 0.0

    async def _variables_for_template(self, template_id: int, template_path: Path) -> list[SmartReportTemplateVariableRow]:
        rows = await self.list_variables(template_id)
        by_key = {item.variable_key: item for item in rows}
        now = _iso_now()
        for idx, token in enumerate(self._extract_placeholders_for_path(template_path)):
            if token in by_key:
                continue
            variable_type, binding_config = self._infer_variable(token)
            rows.append(
                SmartReportTemplateVariableRow(
                    variable_id=-(idx + 1),
                    template_id=template_id,
                    variable_key=token,
                    variable_name=self._default_variable_name(token),
                    variable_type=variable_type,  # type: ignore[arg-type]
                    binding_config=binding_config,
                    display_order=idx,
                    created_at=now,
                    updated_at=now,
                )
            )
        return rows

    def _calc_metric_from_row(self, row: Any) -> SmartReportCalcMetricRow:
        raw_components = json.loads(str(row[3] or "[]"))
        components = [
            SmartReportCalcMetricComponent(
                alias=str(item.get("alias") or ""),
                data_acct_code=str(item.get("data_acct_code") or ""),
                data_acct_name=str(item.get("data_acct_name")) if item.get("data_acct_name") is not None else None,
            )
            for item in raw_components
            if isinstance(item, dict)
        ]
        return SmartReportCalcMetricRow(
            metric_code=str(row[0]),
            metric_name=str(row[1]),
            expression=str(row[2]),
            components=components,
            value_type=str(row[4] or "金额"),
            format_type=str(row[5] or "number"),
            remark=str(row[6]) if row[6] is not None else None,
            created_at=str(row[7]),
            updated_at=str(row[8]),
        )

    def _safe_eval_expression(self, expression: str, values: dict[str, float]) -> float:
        node = ast.parse(expression, mode="eval")

        def visit(current: ast.AST) -> float:
            if isinstance(current, ast.Expression):
                return visit(current.body)
            if isinstance(current, ast.Constant) and isinstance(current.value, (int, float)):
                return float(current.value)
            if isinstance(current, ast.Name):
                if current.id not in values:
                    raise HTTPException(status_code=400, detail=f"表达式引用了未选择的基础指标：{current.id}")
                return float(values[current.id])
            if isinstance(current, ast.BinOp) and type(current.op) in _CALC_OPERATORS:
                return float(_CALC_OPERATORS[type(current.op)](visit(current.left), visit(current.right)))
            if isinstance(current, ast.UnaryOp) and type(current.op) in _CALC_OPERATORS:
                return float(_CALC_OPERATORS[type(current.op)](visit(current.operand)))
            raise HTTPException(status_code=400, detail="表达式只支持基础指标别名、数字和 + - * / ( ) 运算")

        return visit(node)

    async def _render_chart_for_placeholder(self, chart_config_code: str, params: dict[str, Any]) -> bytes:
        """Render a smart-PPT chart config as PNG bytes for Word embedding."""
        if self.smart_ppt_service is None:
            raise HTTPException(status_code=500, detail="智能图表服务未初始化")
        return await self.smart_ppt_service.render_chart_image(chart_config_code, params)

    def _save_chart_cache_png(self, chart_config_code: str, png_bytes: bytes) -> Path:
        self.chart_cache_dir.mkdir(parents=True, exist_ok=True)
        safe_code = re.sub(r"[^A-Za-z0-9_-]+", "_", chart_config_code.strip()).strip("_") or "chart"
        target = self.chart_cache_dir / f"smart_report_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}_{safe_code}.png"
        target.write_bytes(png_bytes)
        return target

    async def _render_docx(
        self,
        template_path: Path,
        output_path: Path,
        resolved: dict[str, str],
        params: dict[str, Any],
    ) -> None:
        if not template_path.exists():
            raise HTTPException(status_code=404, detail="模板文件不存在")
        doc = Document(str(template_path))
        for paragraph in doc.paragraphs:
            await self._replace_docx_paragraph(paragraph, resolved, params)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        await self._replace_docx_paragraph(paragraph, resolved, params)
        doc.save(str(output_path))

    def _render_pptx(self, template_path: Path, output_path: Path, resolved: dict[str, str]) -> None:
        from pptx import Presentation

        if not template_path.exists():
            raise HTTPException(status_code=404, detail="PPTX 模板文件不存在")
        prs = Presentation(str(template_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        self._replace_pptx_paragraph(paragraph, resolved)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                self._replace_pptx_paragraph(paragraph, resolved)
        prs.save(str(output_path))

    def _render_preview_text(
        self,
        template_path: Path,
        resolved: dict[str, str],
        *,
        chart_labels: dict[str, str] | None = None,
    ) -> str:
        if not template_path.exists():
            raise HTTPException(status_code=404, detail="模板文件不存在")
        if template_path.suffix.lower() == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(template_path))
            blocks: list[str] = []
            for index, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = self._replace_preview_text(shape.text_frame.text, resolved, chart_labels or {}).strip()
                        if text:
                            slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [self._replace_preview_text(cell.text, resolved, chart_labels or {}).strip() for cell in row.cells]
                            line = " | ".join(cell for cell in cells if cell)
                            if line:
                                slide_texts.append(line)
                if slide_texts:
                    blocks.append(f"--- Slide {index} ---")
                    blocks.extend(slide_texts)
            return "\n".join(blocks)
        doc = Document(str(template_path))
        blocks: list[str] = []
        for paragraph in doc.paragraphs:
            text = self._replace_preview_text(paragraph.text, resolved, chart_labels or {}).strip()
            if text:
                blocks.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [self._replace_preview_text(cell.text, resolved, chart_labels or {}).strip() for cell in row.cells]
                line = " | ".join(cell for cell in cells if cell)
                if line:
                    blocks.append(line)
        return "\n".join(blocks)

    def _write_text_template_docx(self, target: Path, title: str, content: str) -> None:
        doc = Document()
        if title.strip():
            doc.add_heading(title.strip(), level=1)
        for line in content.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
            doc.add_paragraph(line)
        doc.save(str(target))

    async def _replace_docx_paragraph(self, paragraph: Any, resolved: dict[str, str], params: dict[str, Any]) -> None:
        text = paragraph.text
        if "{{" not in text:
            return
        if not CHART_PLACEHOLDER_RE.search(text):
            self._replace_in_paragraph(paragraph, resolved)
            return

        for run in paragraph.runs:
            run.text = ""
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        cursor = 0
        for match in CHART_PLACEHOLDER_RE.finditer(text):
            before = self._replace_text(text[cursor:match.start()], resolved)
            if before:
                paragraph.add_run(before)

            chart_code = match.group(1).strip()
            png_bytes = await self._render_chart_for_placeholder(chart_code, params)
            chart_path = self._save_chart_cache_png(chart_code, png_bytes)
            paragraph.add_run().add_picture(str(chart_path), width=Inches(6))
            cursor = match.end()

        tail = self._replace_text(text[cursor:], resolved)
        if tail:
            paragraph.add_run(tail)

    def _replace_in_paragraph(self, paragraph: Any, resolved: dict[str, str]) -> None:
        text = paragraph.text
        if "{{" not in text:
            return
        new_text = self._replace_text(text, resolved)
        if not paragraph.runs:
            paragraph.add_run(new_text)
            return
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""

    def _replace_pptx_paragraph(self, paragraph: Any, resolved: dict[str, str]) -> None:
        text = paragraph.text
        if "{{" not in text:
            return
        new_text = self._replace_text(text, resolved)
        if paragraph.runs:
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""
            return
        paragraph.text = new_text

    def _replace_text(self, text: str, resolved: dict[str, str]) -> str:
        return PLACEHOLDER_RE.sub(lambda m: resolved.get(m.group(1).strip(), m.group(0)), text)

    def _replace_preview_text(self, text: str, resolved: dict[str, str], chart_labels: dict[str, str]) -> str:
        def replace_chart(match: re.Match[str]) -> str:
            chart_code = match.group(1).strip()
            return f"[图表: {chart_labels.get(chart_code, chart_code)}]"

        return self._replace_text(CHART_PLACEHOLDER_RE.sub(replace_chart, text), resolved)

    async def _chart_preview_labels(self, template_path: Path) -> dict[str, str]:
        labels: dict[str, str] = {}
        for token in self._extract_placeholders_for_path(template_path):
            prefix, _, raw_code = token.partition(":")
            if prefix.strip().lower() != "chart" or not raw_code.strip():
                continue
            chart_code = raw_code.strip()
            labels[chart_code] = chart_code
            if self.smart_ppt_service is None:
                continue
            chart_config = await self.smart_ppt_service.get_chart_config_by_code(chart_code)
            if not chart_config:
                continue
            visual_title = str(chart_config.visual_config_json.get("title") or "").strip()
            labels[chart_code] = visual_title or str(chart_config.remark or "").strip() or self._chart_type_label(chart_config.chart_type)
        return labels

    def _chart_type_label(self, chart_type: str) -> str:
        return {
            "line": "折线图",
            "bar": "柱状图",
            "dual_bar": "双柱图",
            "donut": "环形图",
            "waterfall": "瀑布图",
        }.get(chart_type, chart_type)

    def _extract_placeholders_for_path(self, template_path: Path) -> list[str]:
        if template_path.suffix.lower() == ".pptx":
            return self.extract_placeholders_pptx(template_path)
        return self.extract_placeholders(template_path)

    def _iter_doc_text(self, doc: Any) -> list[str]:
        texts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    texts.extend(p.text for p in cell.paragraphs)
        return texts
