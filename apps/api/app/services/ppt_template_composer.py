"""Compose PPTX files from a source template and binding configuration."""

from __future__ import annotations

from dataclasses import dataclass
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.schemas import SmartPptTemplateBindingConfigRow


@dataclass
class PptTemplateComposeResult:
    applied_count: int
    slide_count: int
    warnings: list[str]


class PptTemplateComposer:
    """Apply lightweight binding updates while preserving template styling."""

    def compose(
        self,
        *,
        template_path: Path,
        output_path: Path,
        bindings: list[SmartPptTemplateBindingConfigRow],
        params: dict[str, Any] | None = None,
        chart_payloads: dict[str, dict[str, Any]] | None = None,
        max_slides: int | None = None,
    ) -> PptTemplateComposeResult:
        prs = Presentation(str(template_path))
        params = params or {}
        chart_payloads = chart_payloads or {}
        warnings: list[str] = []
        applied_count = 0

        for binding in bindings:
            if not binding.enabled or binding.binding_type == "ignore":
                continue
            shape = self._shape_by_object_id(prs, binding.object_id)
            if shape is None:
                warnings.append(f"{binding.object_id}：未在模板中找到对应对象。")
                continue

            if binding.binding_type in {"text", "kpi"}:
                if not getattr(shape, "has_text_frame", False):
                    warnings.append(f"{binding.object_id}：不是文本对象，已跳过文本绑定。")
                    continue
                text_value = self._binding_text(binding, params)
                if text_value is None:
                    warnings.append(f"{binding.object_id}：没有提供 {binding.target_key or 'target_key'} 的文本值，已保留模板原文。")
                    continue
                self._write_text(shape, text_value)
                applied_count += 1
            elif binding.binding_type == "table":
                if not getattr(shape, "has_table", False):
                    warnings.append(f"{binding.object_id}：不是表格对象，已跳过表格绑定。")
                    continue
                table_payload = self._table_payload(binding, params)
                if table_payload is None:
                    warnings.append(f"{binding.object_id}：没有提供 {binding.target_key or 'target_key'} 的表格数据，已保留模板原表格。")
                    continue
                self._write_table(shape.table, table_payload)
                applied_count += 1
            elif binding.binding_type == "chart":
                if not getattr(shape, "has_chart", False):
                    warnings.append(f"{binding.object_id}：不是图表对象，已跳过图表绑定。")
                    continue
                chart_payload = chart_payloads.get(binding.object_id)
                if not chart_payload:
                    warnings.append(f"{binding.object_id}：缺少图表数据，已保留模板原样。")
                    continue
                self._replace_chart_data(shape.chart, chart_payload)
                applied_count += 1
            else:
                warnings.append(f"{binding.object_id}：暂不支持绑定类型 {binding.binding_type}。")

        if applied_count == 0:
            warnings.append("没有可应用的绑定，已输出原模板副本。")

        if max_slides is not None:
            self._trim_slides(prs, max_slides)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return PptTemplateComposeResult(applied_count=applied_count, slide_count=len(prs.slides), warnings=warnings)

    def _shape_by_object_id(self, prs: Presentation, object_id: str) -> Any | None:
        try:
            slide_token, path = object_id.split(":", 1)
            slide_index = int(slide_token.removeprefix("s"))
        except Exception:
            return None
        if slide_index < 1 or slide_index > len(prs.slides):
            return None

        target_path = path.strip()
        for current_path, shape in self._iter_shapes(prs.slides[slide_index - 1].shapes):
            if current_path == target_path:
                return shape
        return None

    def _iter_shapes(self, shapes: Iterable[Any], prefix: str = "") -> Iterable[tuple[str, Any]]:
        for index, shape in enumerate(shapes, start=1):
            path = f"{prefix}.{index}" if prefix else str(index)
            yield path, shape
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                yield from self._iter_shapes(shape.shapes, path)

    def _trim_slides(self, prs: Presentation, max_slides: int) -> None:
        if max_slides < 1:
            return
        slide_id_list = prs.slides._sldIdLst  # python-pptx has no public delete API.
        while len(prs.slides) > max_slides:
            slide_id = slide_id_list[-1]
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            slide_id_list.remove(slide_id)

    def _binding_text(self, binding: SmartPptTemplateBindingConfigRow, params: dict[str, Any]) -> str | None:
        for key in [binding.target_key]:
            if key and key in params:
                return str(params[key])
        return None

    def _write_text(self, shape: Any, value: str) -> None:
        frame = shape.text_frame
        paragraph = frame.paragraphs[0]
        if paragraph.runs:
            paragraph.runs[0].text = value
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = value
        for extra_paragraph in frame.paragraphs[1:]:
            for run in extra_paragraph.runs:
                run.text = ""

    def _table_payload(self, binding: SmartPptTemplateBindingConfigRow, params: dict[str, Any]) -> dict[str, Any] | None:
        if not binding.target_key or binding.target_key not in params:
            return None
        value = params[binding.target_key]
        if isinstance(value, dict):
            headers = value.get("headers") or []
            rows = value.get("rows") or []
            return {"headers": headers, "rows": rows}
        if isinstance(value, list):
            return {"headers": [], "rows": value}
        return None

    def _write_table(self, table: Any, payload: dict[str, Any]) -> None:
        rows = list(table.rows)
        columns = list(table.columns)
        if not rows or not columns:
            return

        headers = [str(item) for item in payload.get("headers") or []]
        data_rows = payload.get("rows") or []
        for col_index, header in enumerate(headers[: len(columns)]):
            rows[0].cells[col_index].text = header
        start_row = 1 if headers else 0
        for row_offset, row_values in enumerate(list(data_rows)[: max(0, len(rows) - start_row)]):
            row_index = start_row + row_offset
            if not isinstance(row_values, list):
                row_values = [row_values]
            for col_index, value in enumerate(row_values[: len(columns)]):
                rows[row_index].cells[col_index].text = str(value)

    def _replace_chart_data(self, chart: Any, payload: dict[str, Any]) -> None:
        labels = [str(label) for label in payload.get("labels", [])]
        if not labels:
            labels = ["-"]

        chart_data = CategoryChartData()
        chart_data.categories = labels

        series_rows = payload.get("series", [])
        if not series_rows:
            series_rows = [{"name": payload.get("series_name") or "数据", "values": payload.get("values") or [0]}]

        for series in series_rows:
            name = str(series.get("name") or "数据")
            values = self._normalize_series_values(series.get("values"), len(labels))
            chart_data.add_series(name, values)

        chart.replace_data(chart_data)

    def _normalize_series_values(self, values: Any, length: int) -> list[float]:
        result: list[float] = []
        for value in list(values or [])[:length]:
            try:
                result.append(float(value or 0))
            except Exception:
                result.append(0.0)
        while len(result) < length:
            result.append(0.0)
        return result
