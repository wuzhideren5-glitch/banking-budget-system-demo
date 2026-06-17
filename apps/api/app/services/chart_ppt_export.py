"""PowerPoint export builder for multidimensional chart results."""
from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone
from io import BytesIO

from fastapi import HTTPException
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas import ChartPptExportRequestDto


PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@dataclass(frozen=True)
class ChartPptExportFile:
    content: BytesIO
    filename: str
    media_type: str = PPTX_MEDIA_TYPE


def _validate_chart_ppt_request(req: ChartPptExportRequestDto) -> list[str]:
    categories = [str(c) for c in req.categories]
    if not categories:
        raise HTTPException(status_code=400, detail="categories 不能为空")
    if not req.series:
        raise HTTPException(status_code=400, detail="series 不能为空")

    for series in req.series:
        if len(series.values) != len(categories):
            raise HTTPException(
                status_code=400,
                detail=f"系列 {series.name} 的数据点数量与 categories 不一致",
            )
    return categories


def build_chart_ppt_export_file(req: ChartPptExportRequestDto) -> ChartPptExportFile:
    categories = _validate_chart_ppt_request(req)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    emu_per_inch = 914400
    slide_width_in = prs.slide_width / emu_per_inch
    slide_height_in = prs.slide_height / emu_per_inch
    left_margin_in = 0.45
    right_margin_in = 0.45
    top_margin_in = 0.18
    bottom_margin_in = 0.18
    content_width_in = slide_width_in - left_margin_in - right_margin_in

    title_box = slide.shapes.add_textbox(
        Inches(left_margin_in), Inches(top_margin_in), Inches(content_width_in), Inches(0.40)
    )
    title_box.text_frame.clear()
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = req.title
    title_para.font.size = Pt(16)
    title_para.font.bold = True

    subtitle_top_in = top_margin_in + 0.40 + 0.04
    chart_top_in = subtitle_top_in
    if req.subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(left_margin_in), Inches(subtitle_top_in), Inches(content_width_in), Inches(0.24)
        )
        subtitle_box.text_frame.clear()
        subtitle_para = subtitle_box.text_frame.paragraphs[0]
        subtitle_para.text = req.subtitle
        subtitle_para.font.size = Pt(9)
        chart_top_in = subtitle_top_in + 0.24 + 0.08

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series in req.series:
        chart_data.add_series(series.name, series.values)

    chart_type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    chart_type = chart_type_map[req.chart_type]

    has_matrix = bool(req.matrix_headers) and bool(req.matrix_rows)
    matrix_top_in = 0.0
    matrix_height_in = 0.0
    matrix_note_top_in = 0.0
    matrix_note_height_in = 0.0
    if has_matrix:
        matrix_note_height_in = 0.16
        matrix_height_in = min(2.40, slide_height_in * 0.33)
        matrix_note_top_in = slide_height_in - bottom_margin_in - matrix_note_height_in
        matrix_top_in = matrix_note_top_in - 0.03 - matrix_height_in
        chart_bottom_limit_in = matrix_top_in - 0.10
    else:
        chart_bottom_limit_in = slide_height_in - bottom_margin_in
    chart_height = max(2.60, chart_bottom_limit_in - chart_top_in)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(left_margin_in),
        Inches(chart_top_in),
        Inches(content_width_in),
        Inches(chart_height),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.has_title = False

    if req.chart_type in {"pie", "doughnut"}:
        plot = chart.plots[0]
        plot.has_data_labels = True
        labels = plot.data_labels
        is_doughnut = req.chart_type == "doughnut"
        labels.show_category_name = True
        labels.show_percentage = not is_doughnut
        labels.show_value = is_doughnut
        labels.number_format = "0.00" if is_doughnut else "0.00%"
        labels.position = XL_LABEL_POSITION.BEST_FIT
        if hasattr(labels, "separator"):
            labels.separator = "\n"
        if hasattr(labels, "font"):
            labels.font.size = Pt(7)
        for chart_series in chart.series:
            for point in chart_series.points:
                point_label = point.data_label
                point_label.show_percentage = not is_doughnut
                point_label.show_value = is_doughnut
                point_label.show_category_name = True
                point_label.number_format = "0.00" if is_doughnut else "0.00%"
                point_label.position = XL_LABEL_POSITION.BEST_FIT
                point_label.font.size = Pt(7)
    else:
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(7)
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(7)

    for plot in chart.plots:
        if getattr(plot, "has_data_labels", False):
            data_labels = plot.data_labels
            if hasattr(data_labels, "font"):
                data_labels.font.size = Pt(7)

    if has_matrix:
        max_data_columns = 8
        max_data_rows = 10

        matrix_headers = [str(h) for h in req.matrix_headers]
        matrix_rows: list[tuple[str, list[str]]] = []
        for row in req.matrix_rows:
            values = [str(v) for v in row.values]
            if len(values) < len(matrix_headers):
                values = values + [""] * (len(matrix_headers) - len(values))
            elif len(values) > len(matrix_headers):
                values = values[: len(matrix_headers)]
            matrix_rows.append((str(row.label), values))

        clipped_columns = len(matrix_headers) > max_data_columns
        clipped_rows = len(matrix_rows) > max_data_rows
        display_headers = matrix_headers[:max_data_columns]
        display_rows = matrix_rows[:max_data_rows]

        table_rows = 1 + len(display_rows)
        table_cols = 1 + len(display_headers)
        table_shape = slide.shapes.add_table(
            table_rows,
            table_cols,
            Inches(left_margin_in),
            Inches(matrix_top_in),
            Inches(content_width_in),
            Inches(matrix_height_in),
        )
        table = table_shape.table
        first_col_width = min(4.20, max(3.20, content_width_in * 0.34))
        table.columns[0].width = Inches(first_col_width)
        value_col_width = Inches(max(0.60, (content_width_in - first_col_width) / max(1, len(display_headers))))
        for col_idx in range(1, table_cols):
            table.columns[col_idx].width = value_col_width

        head_cell = table.cell(0, 0)
        head_cell.text = "维度项"
        head_para = head_cell.text_frame.paragraphs[0]
        head_para.font.bold = True
        head_para.font.size = Pt(11)
        head_para.alignment = PP_ALIGN.CENTER

        for col_idx, header in enumerate(display_headers, start=1):
            cell = table.cell(0, col_idx)
            cell.text = header
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER

        for row_idx, (row_label, values) in enumerate(display_rows, start=1):
            label_cell = table.cell(row_idx, 0)
            label_cell.text = row_label
            label_para = label_cell.text_frame.paragraphs[0]
            label_para.font.size = Pt(10)
            label_para.alignment = PP_ALIGN.LEFT

            for col_idx, value in enumerate(values[: len(display_headers)], start=1):
                value_cell = table.cell(row_idx, col_idx)
                value_cell.text = value
                value_para = value_cell.text_frame.paragraphs[0]
                value_para.font.size = Pt(10)
                value_para.alignment = PP_ALIGN.RIGHT

        if clipped_columns or clipped_rows:
            note_box = slide.shapes.add_textbox(
                Inches(left_margin_in), Inches(matrix_note_top_in), Inches(content_width_in), Inches(matrix_note_height_in)
            )
            note_box.text_frame.text = "注：数据矩阵已按页面宽度/高度截取显示。"
            note_box.text_frame.paragraphs[0].font.size = Pt(7)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    filename = f"pivot_chart_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.pptx"
    return ChartPptExportFile(content=output, filename=filename)
