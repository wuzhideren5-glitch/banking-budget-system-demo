"""Editable PowerPoint rendering for scene-driven smart PPT decks."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas import SmartPptChartConfigRow, SmartPptSceneRow

SUPPORTED_NATIVE_CHART_TYPES = {"line", "bar", "dual_bar", "donut"}


class SmartPptRenderer:
    """Render slide payloads into editable PPTX content.

    The service builds semantic payloads; this class owns visual style, native
    PowerPoint chart data, slide geometry, and table formatting.
    """

    def compose(self, scene: SmartPptSceneRow, slides: list[dict[str, Any]], output_path: Path) -> None:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for payload in slides:
            slide_spec = payload["slide_spec"]
            slide_type = str(slide_spec.get("type") or "text")
            title = payload["title"]
            subtitle = payload["subtitle"]
            narrative = payload["narrative"]
            chart_config = payload.get("chart_config")
            metrics = payload["metrics"]

            if slide_type == "cover":
                self._add_cover_slide(prs, title, subtitle or scene.scene_name)
            elif slide_type == "dashboard":
                self._add_dashboard_slide(prs, title, narrative, metrics.get("metric_cards", []))
            elif chart_config is not None and chart_config.chart_type in SUPPORTED_NATIVE_CHART_TYPES:
                self._add_chart_slide(
                    prs,
                    title,
                    narrative,
                    chart_config,
                    metrics,
                    metrics.get("table_headers", []),
                    metrics.get("table_rows", []),
                )
            else:
                self._add_text_slide(prs, title, narrative)

        prs.save(str(output_path))

    def _add_cover_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x0F, 0x3D, 0x5C)
        accent.line.fill.background()

        kicker_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(10.8), Inches(0.32))
        kicker = kicker_box.text_frame.paragraphs[0]
        kicker.text = "BANKING BUDGET INTELLIGENCE"
        kicker.font.size = Pt(10)
        kicker.font.bold = True
        kicker.font.color.rgb = RGBColor(0x2F, 0x66, 0x75)

        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(10.9), Inches(1.6))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title
        title_paragraph.font.size = Pt(34)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(0x17, 0x24, 0x33)

        subtitle_box = slide.shapes.add_textbox(Inches(0.92), Inches(4.0), Inches(10.8), Inches(0.55))
        subtitle_paragraph = subtitle_box.text_frame.paragraphs[0]
        subtitle_paragraph.text = subtitle
        subtitle_paragraph.font.size = Pt(16)
        subtitle_paragraph.font.color.rgb = RGBColor(0x5A, 0x6B, 0x7C)

        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(5.05), Inches(2.0), Inches(0.04))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(0xC8, 0x8A, 0x3D)
        rule.line.fill.background()

        footer = slide.shapes.add_textbox(Inches(0.92), Inches(6.7), Inches(7.2), Inches(0.28))
        footer_para = footer.text_frame.paragraphs[0]
        footer_para.text = "Generated from budget summary data"
        footer_para.font.size = Pt(8)
        footer_para.font.color.rgb = RGBColor(0x7B, 0x88, 0x94)

    def _add_dashboard_slide(self, prs: Presentation, title: str, narrative: str, metric_cards: list[dict[str, str]]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_title_bar(slide, title)

        narrative_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.28), Inches(11.9), Inches(0.72))
        narrative_panel.fill.solid()
        narrative_panel.fill.fore_color.rgb = RGBColor(0xF3, 0xF6, 0xF8)
        narrative_panel.line.color.rgb = RGBColor(0xDE, 0xE6, 0xEA)

        narrative_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.43), Inches(11.35), Inches(0.42))
        narrative_frame = narrative_box.text_frame
        narrative_frame.word_wrap = True
        paragraph = narrative_frame.paragraphs[0]
        paragraph.text = narrative
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = RGBColor(0x24, 0x34, 0x46)

        for index, card_data in enumerate(metric_cards[:4]):
            left = Inches(0.72 + (index % 2) * 6.18)
            top = Inches(2.35 + (index // 2) * 1.82)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.72), Inches(1.42))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            card.line.color.rgb = RGBColor(0xD8, 0xE2, 0xE8)

            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), Inches(1.42))
            accent.fill.solid()
            accent.fill.fore_color.rgb = next(self._card_colors())
            accent.line.fill.background()

            frame = card.text_frame
            frame.clear()
            frame.margin_left = Inches(0.22)
            frame.margin_right = Inches(0.16)
            frame.margin_top = Inches(0.12)
            frame.margin_bottom = Inches(0.08)
            header = frame.paragraphs[0]
            header.text = str(card_data.get("指标") or "")
            header.font.size = Pt(10)
            header.font.bold = True
            header.font.color.rgb = RGBColor(0x5C, 0x6B, 0x78)

            value = frame.add_paragraph()
            value.text = str(card_data.get("实际") or "-")
            value.font.size = Pt(22)
            value.font.bold = True
            value.font.color.rgb = RGBColor(0x17, 0x24, 0x33)

            details = frame.add_paragraph()
            details.text = f"预算 {card_data.get('预算', '-')}  差异 {card_data.get('差异', '-')}"
            details.font.size = Pt(8.5)
            details.font.color.rgb = RGBColor(0x58, 0x66, 0x73)

            footer = frame.add_paragraph()
            footer.text = f"完成率 {card_data.get('完成率', '-')}  同比 {card_data.get('同比', '-')}"
            footer.font.size = Pt(8.5)
            footer.font.color.rgb = RGBColor(0x58, 0x66, 0x73)

    def _add_chart_slide(
        self,
        prs: Presentation,
        title: str,
        narrative: str,
        chart_config: SmartPptChartConfigRow,
        metrics: dict[str, Any],
        table_headers: list[str],
        table_rows: list[list[str]],
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_title_bar(slide, title)

        narrative_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.25), Inches(11.65), Inches(0.62))
        narrative_frame = narrative_box.text_frame
        narrative_frame.word_wrap = True
        narrative_paragraph = narrative_frame.paragraphs[0]
        narrative_paragraph.text = narrative
        narrative_paragraph.font.size = Pt(12.5)
        narrative_paragraph.font.color.rgb = RGBColor(0x37, 0x47, 0x58)

        chart_data = self._build_native_chart_data(chart_config, metrics)
        chart_shape = slide.shapes.add_chart(
            self._native_chart_type(chart_config.chart_type),
            Inches(0.68),
            Inches(2.0),
            Inches(8.0),
            Inches(4.72),
            chart_data,
        )
        self._format_native_chart(chart_shape.chart, chart_config.chart_type)
        if table_headers and table_rows:
            self._add_summary_table(slide, table_headers, table_rows[:6], Inches(8.92), Inches(2.0), Inches(3.72), Inches(4.72))

    def _native_chart_type(self, chart_type: str) -> Any:
        mapping = {
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "dual_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "donut": XL_CHART_TYPE.DOUGHNUT,
        }
        return mapping[chart_type]

    def _build_native_chart_data(self, chart_config: SmartPptChartConfigRow, metrics: dict[str, Any]) -> CategoryChartData:
        labels = [str(label) for label in metrics.get("labels") or []]
        if not labels:
            labels = ["暂无数据"]

        chart_type = chart_config.chart_type
        chart_data = CategoryChartData()
        chart_data.categories = labels

        def values_of(key: str, *, fallback_key: str | None = None) -> list[float]:
            raw_values = metrics.get(key)
            if raw_values is None and fallback_key:
                raw_values = metrics.get(fallback_key)
            values = [float(value or 0) for value in (raw_values or [])]
            if len(values) < len(labels):
                values = values + [0.0] * (len(labels) - len(values))
            return values[: len(labels)]

        metric_name = str(metrics.get("metric_name") or chart_config.visual_config_json.get("title") or "指标")
        if chart_type == "line":
            chart_data.add_series("实际", values_of("values"))
            budget_values = values_of("budget_values")
            if any(abs(value) > 1e-9 for value in budget_values):
                chart_data.add_series("预算", budget_values)
        elif chart_type == "dual_bar":
            chart_data.add_series("实际", values_of("actual_values", fallback_key="values"))
            chart_data.add_series("预算", values_of("budget_values"))
        elif chart_type == "bar":
            chart_data.add_series("实际", values_of("values"))
            budget_values = values_of("budget_values")
            if any(abs(value) > 1e-9 for value in budget_values):
                chart_data.add_series("预算", budget_values)
        elif chart_type == "donut":
            values = values_of("values")
            if not any(abs(value) > 1e-9 for value in values):
                chart_data = CategoryChartData()
                chart_data.categories = ["暂无数据"]
                values = [1.0]
            chart_data.add_series(metric_name, values)
        return chart_data

    def _format_native_chart(self, chart: Any, chart_type: str) -> None:
        chart.has_title = False
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
        self._apply_chart_palette(chart, chart_type)

        if chart_type == "donut":
            plot = chart.plots[0]
            plot.has_data_labels = True
            labels = plot.data_labels
            labels.show_category_name = True
            labels.show_percentage = True
            labels.show_value = False
            labels.number_format = "0.0%"
            labels.position = XL_LABEL_POSITION.BEST_FIT
            if hasattr(labels, "separator"):
                labels.separator = "\n"
            labels.font.size = Pt(7)
            for series in chart.series:
                for point in series.points:
                    point_label = point.data_label
                    point_label.show_category_name = True
                    point_label.show_percentage = True
                    point_label.show_value = False
                    point_label.number_format = "0.0%"
                    point_label.position = XL_LABEL_POSITION.BEST_FIT
                    point_label.font.size = Pt(7)
            return

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(8)
        value_axis = chart.value_axis
        value_axis.tick_labels.font.size = Pt(8)
        value_axis.has_major_gridlines = True

    def _apply_chart_palette(self, chart: Any, chart_type: str) -> None:
        colors = [
            RGBColor(0x1B, 0x5E, 0x7A),
            RGBColor(0xC8, 0x8A, 0x3D),
            RGBColor(0x5F, 0x7F, 0x52),
            RGBColor(0x8B, 0x5E, 0x83),
        ]
        for index, series in enumerate(chart.series):
            color = colors[index % len(colors)]
            try:
                if chart_type == "line":
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(2)
                else:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
            except Exception:
                pass

    def _add_summary_table(
        self,
        slide: Any,
        headers: list[str],
        rows: list[list[str]],
        left: Any,
        top: Any,
        width: Any,
        height: Any,
    ) -> None:
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
        table = table_shape.table
        if len(headers) > 0:
            table.columns[0].width = Inches(1.05)
        if len(headers) > 1:
            remaining_width = max(0.6, (float(width) / 914400 - 1.05) / (len(headers) - 1))
            for col_index in range(1, len(headers)):
                table.columns[col_index].width = Inches(remaining_width)
        for col_index, header in enumerate(headers):
            cell = table.cell(0, col_index)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x21, 0x32, 0x43)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = Pt(8.5)
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
        for row_index, row in enumerate(rows, start=1):
            for col_index, value in enumerate(row[: len(headers)]):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if row_index % 2 else RGBColor(0xF5, 0xF7, 0xF9)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(7.8)
                paragraph.font.color.rgb = RGBColor(0x2F, 0x3D, 0x4A)
                paragraph.alignment = PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.RIGHT

    def _add_text_slide(self, prs: Presentation, title: str, narrative: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_title_bar(slide, title)
        text_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.2), Inches(4.8))
        frame = text_box.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = narrative
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(0x24, 0x34, 0x46)

    def _add_title_bar(self, slide: Any, title: str) -> None:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(0.62), Inches(0.11), Inches(0.38))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0xC8, 0x8A, 0x3D)
        accent.line.fill.background()

        box = slide.shapes.add_textbox(Inches(0.84), Inches(0.48), Inches(10.9), Inches(0.52))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x17, 0x24, 0x33)

        label = slide.shapes.add_textbox(Inches(10.7), Inches(0.56), Inches(1.95), Inches(0.3))
        label_para = label.text_frame.paragraphs[0]
        label_para.text = "Smart PPT"
        label_para.font.size = Pt(8)
        label_para.font.color.rgb = RGBColor(0x7B, 0x88, 0x94)
        label_para.alignment = PP_ALIGN.RIGHT

        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.66), Inches(1.12), Inches(11.98), Inches(0.01))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(0xD7, 0xE0, 0xE5)
        rule.line.fill.background()

    def _set_slide_background(self, slide: Any) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFD)

    @staticmethod
    def _card_colors():
        colors = [
            RGBColor(0x1B, 0x5E, 0x7A),
            RGBColor(0xC8, 0x8A, 0x3D),
            RGBColor(0x5F, 0x7F, 0x52),
            RGBColor(0x8B, 0x5E, 0x83),
        ]
        while True:
            for color in colors:
                yield color
