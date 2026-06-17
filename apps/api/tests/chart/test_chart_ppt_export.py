from __future__ import annotations

from io import BytesIO
import unittest

from fastapi import HTTPException
from pptx import Presentation

from app.schemas import ChartPptExportRequestDto, ChartPptMatrixRowDto, ChartPptSeriesDto
from app.services.chart_ppt_export import PPTX_MEDIA_TYPE, build_chart_ppt_export_file


class ChartPptExportTests(unittest.TestCase):
    def test_builds_pptx_with_chart_matrix_and_clipping_note(self) -> None:
        req = ChartPptExportRequestDto(
            chart_type="bar",
            title="经营分析",
            subtitle="多年度数据透视图",
            categories=[f"C{i}" for i in range(1, 10)],
            series=[ChartPptSeriesDto(name="余额", values=[float(i) for i in range(1, 10)])],
            matrix_headers=[f"H{i}" for i in range(1, 10)],
            matrix_rows=[
                ChartPptMatrixRowDto(label=f"R{i}", values=[str(i)] * 9)
                for i in range(1, 12)
            ],
        )

        export_file = build_chart_ppt_export_file(req)
        content = export_file.content.getvalue()

        self.assertEqual(export_file.media_type, PPTX_MEDIA_TYPE)
        self.assertRegex(export_file.filename, r"^pivot_chart_\d{8}_\d{6}\.pptx$")
        self.assertTrue(content.startswith(b"PK"))

        prs = Presentation(BytesIO(content))
        self.assertEqual(len(prs.slides), 1)
        self.assertGreater(prs.slide_width, prs.slide_height)

        texts = []
        for shape in prs.slides[0].shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
        joined = "\n".join(texts)
        self.assertIn("经营分析", joined)
        self.assertIn("多年度数据透视图", joined)
        self.assertIn("注：数据矩阵已按页面宽度/高度截取显示。", joined)

    def test_rejects_series_length_mismatch(self) -> None:
        req = ChartPptExportRequestDto(
            chart_type="line",
            title="经营分析",
            categories=["M01", "M02"],
            series=[ChartPptSeriesDto(name="余额", values=[1.0])],
        )

        with self.assertRaises(HTTPException) as ctx:
            build_chart_ppt_export_file(req)

        self.assertEqual(ctx.exception.status_code, 400)
        self.assertIn("数据点数量", str(ctx.exception.detail))


if __name__ == "__main__":
    unittest.main()
